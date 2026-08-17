#!/usr/bin/env python3
"""COMPACT v2 slide-deck reference — the information-dense visual format.

This is the canonical reference for the Tertiary Infotech COMPACT deck format
(origin: the SPC in Manufacturing course TGS-2026064862, v14 — kept here as a
worked example, exactly as the AZ-104 pipeline is kept for the classic format).
Copy it into a new course repo and swap the content; the component library is
course-agnostic and everything is driven by course_data.py + data_domainN.py.

WHAT V2 ADDS over the classic reference (build_slides.py):

  img_points(title, image, points)   chart/diagram image + takeaway tiles — the
                                     workhorse concept slide (image left, 3-4
                                     colour-coded takeaways right)
  img_full(title, image, caption)    full-width visual with a caption band
  table_slide(title, headers, rows)  styled decision/comparison tables
  formula_slide(title, panels)       formula panels (dark code-style boxes)
  activity_slide(a, topic)           ONE workflow slide per hands-on activity:
                                     badge + scenario + 5-chip numbered strip +
                                     "YOU'LL PRODUCE" band + Tools line.
                                     NEVER one-step-per-slide runs.
  lms_slide() / novaspc_slide()      browser-mockup visual slides (URL bar,
                                     traffic lights, nav mock, numbered how-to
                                     tiles) for the LMS portal and any course
                                     ed-tool — never a bare text link
  slide_map.json export              mark()/SLIDE_MAP records where each topic
                                     and activity starts so the Lesson Plan can
                                     cite exact slide numbers and never drift

DATA CONTRACT ADDITIONS (course_data.py + data_domainN.py):
  - TOPICS[n]["concepts"] entries are (title, caption) TUPLES → tile_grid renders
    bold title + grey caption (much denser than plain strings)
  - each activity dict adds: minutes=N, flow=[5 short chip labels], and
    optionally csv=dict(name=..., rows=[[header],[...]]) written to labs/data/
  - course ed-tool (if any) goes in services + a flow chip ("Upload x.csv to
    <tool>" → "<panel> → Generate → compare") — hand-compute first, tool second

CHART ASSETS: generate with make_charts.py (sibling reference) — matplotlib,
Arial, white bg, brand palette, 150 dpi. Draw every chart FROM THE SAME NUMBERS
the activities use so slides, LG and assessment can never disagree.

HARD-LEARNED RULES (each cost a QA round somewhere):
  - slide titles <= ~48 chars or they wrap into the divider rule
  - fig.tight_layout(rect=[0,0,1,0.86]) when a suptitle exists; h_pad>=2 for
    stacked subplots; white bbox behind any label that sits ON a chart line
  - stick to glyphs Arial has (no superscript minus U+207B)
  - every generated asset must actually be placed on a slide
  - image placement: aspect-fit with PIL against the content area, never stretch
"""

import os, sys, json, math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
from data_domain1 import DOMAIN1
from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3
from data_domain4 import DOMAIN4
from data_domain5 import DOMAIN5
ACTIVITIES = DOMAIN1 + DOMAIN2 + DOMAIN3 + DOMAIN4 + DOMAIN5

def _find_repo(start):
    env = os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env): return env
    d = start
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "labs")): return d
    return os.path.dirname(os.path.dirname(HERE))
REPO = _find_repo(HERE)
ASSETS = os.path.join(REPO, "courseware", "assets")

# ---------------- palette ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)
RED=RGBColor(0xDC,0x26,0x26); NAVY=RGBColor(0x0B,0x12,0x20)
PALETTE=[BLUE,TEAL,VIOLET,AMBER]

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=color if lvl==0 else GREY
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

PAGE={"n":0}; SLIDE_MAP={}
def mark(key):
    """Record that <key> begins on the NEXT slide to be numbered."""
    SLIDE_MAP[key]=PAGE["n"]+1
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}",9,GREY,False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,GREY,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,GREY,False)]],align=PP_ALIGN.RIGHT)
def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),Inches(1.55),kcolor)
    if kicker: txt(s,Inches(0.85),Inches(0.5),Inches(11.6),Inches(0.4),[[(kicker,14,kcolor,True)]])
    txt(s,Inches(0.85),Inches(0.9),Inches(11.9),Inches(0.9),[[(title,28,INK,True)]])
    rect(s,Inches(0.85),Inches(1.7),Inches(11.63),Inches(0.02),LINE)
    return s
def _asset(name):
    p=os.path.join(ASSETS,name)
    return p if os.path.exists(p) else None

# ---------------- component library (reference deck) ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_asset("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    # course badge (top-right)
    rect(s,Inches(10.85),Inches(0.72),Inches(1.7),Inches(1.0),BLUE)
    txt(s,Inches(10.85),Inches(0.84),Inches(1.7),Inches(0.5),[[("SPC",22,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(10.85),Inches(1.32),Inches(1.7),Inches(0.4),[[("QUALITY CONTROL",8,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,38,INK,True)]])
    rect(s,Inches(0.92),Inches(4.45),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(4.8),Inches(12),Inches(1.5),
        [[(f"WSQ Course Code: {C.COURSE_CODE}  ·  1-Day Hands-On Training",16,GREY,False)],
         [(f"Skills Framework TSC: {C.TSC_TITLE} ({C.TSC_CODE})",14,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.45),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])
    PAGE["n"]+=1

def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,LINE,True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    s=head(slide(),title,kicker); bullets(s,Inches(0.85),Inches(1.95),Inches(11.6),Inches(4.9),items,size=size); footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead="",lcolor=BLUE,rcolor=TEAL,note=None):
    s=head(slide(),title,kicker)
    bh=Inches(4.7) if not note else Inches(4.15)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),bh,LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),bh,LIGHT)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(0.1),lcolor); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(0.1),rcolor)
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,lcolor,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,rcolor,True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),bh-Inches(0.9),left,size=15)
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),bh-Inches(0.9),right,size=15,mcolor=rcolor)
    if note:
        txt(s,Inches(0.85),Inches(6.25),Inches(11.7),Inches(0.6),[[(note,13,GREY,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def cards3(title,cards,kicker):
    s=head(slide(),title,kicker); xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),Inches(3.65),Inches(4.7),LIGHT); rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),Inches(3.2),Inches(3.4),c[2],size=14,mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,INK,True)],[(it[1],size-2,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def flow_h(title,steps,kicker=None,color=BLUE,note=None):
    s=head(slide(),title,kicker,kcolor=color)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.45); ch=Inches(3.15); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.16),int(y+Inches(1.55)),cw-Inches(0.32),int(ch-Inches(1.7)),[[(st,13,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if note:
        txt(s,Inches(0.85),Inches(5.95),Inches(11.7),Inches(0.7),[[(note,14,GREY,False)]],align=PP_ALIGN.CENTER)
    footer(s); return s
def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),LIGHT); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,GREY,False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,INK,False)] if val else [("____________________________________________",13,LINE,False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1

# ---------------- information-dense additions (reference-deck style) ----------------
def img_points(title,image,points,kicker=None,accent=BLUE,img_w=7.0,note=None):
    """Chart image on the left, takeaway tiles on the right — the compact
    concept layout used throughout the reference deck."""
    s=head(slide(),title,kicker,kcolor=accent)
    p=_asset(image)
    area_h=Inches(4.75) if not note else Inches(4.35)
    if p:
        from PIL import Image
        iw,ih=Image.open(p).size
        w=Inches(img_w); h=int(w*ih/iw)
        if h>area_h: h=area_h; w=int(h*iw/ih)
        y=int(Inches(1.95)+(area_h-h)/2)
        s.shapes.add_picture(p,Inches(0.85),y,width=w,height=h)
        rx=int(Inches(0.85)+w+Inches(0.35))
    else:
        rx=Inches(6.5)
    rw=int(Inches(12.48)-rx)
    n=len(points); gy=Inches(0.2); th=int((area_h-gy*(n-1))/n)
    for i,(pt_title,pt_sub) in enumerate(points):
        y=int(Inches(1.95)+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.09),th,col)
        txt(s,rx+Inches(0.28),y,rw-Inches(0.5),th,
            [[(pt_title,14,col,True)],[(pt_sub,12,INK,False)]],anchor=MSO_ANCHOR.MIDDLE,space=2)
    if note:
        txt(s,Inches(0.85),Inches(6.4),Inches(11.7),Inches(0.5),[[(note,13,GREY,False)]],align=PP_ALIGN.CENTER)
    footer(s); return s
def img_full(title,image,kicker=None,accent=BLUE,caption=None):
    """Full-width chart image with an optional caption band."""
    s=head(slide(),title,kicker,kcolor=accent)
    p=_asset(image)
    area_h=Inches(4.35) if caption else Inches(4.8)
    if p:
        from PIL import Image
        iw,ih=Image.open(p).size
        w=Inches(11.6); h=int(w*ih/iw)
        if h>area_h: h=area_h; w=int(h*iw/ih)
        x=int((SW-w)/2); y=int(Inches(1.9)+(area_h-h)/2)
        s.shapes.add_picture(p,x,y,width=w,height=h)
    if caption:
        rect(s,Inches(0.85),Inches(6.3),Inches(11.63),Inches(0.62),LIGHT)
        txt(s,Inches(1.1),Inches(6.3),Inches(11.2),Inches(0.62),[[(caption,14,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def table_slide(title,headers,rows,kicker=None,accent=BLUE,widths=None,note=None,fsize=13):
    """Styled data table — comparison/decision tables like the reference deck."""
    s=head(slide(),title,kicker,kcolor=accent)
    ncol=len(headers); X0=Inches(0.85); TOTW=Inches(11.63)
    if widths: ws=[int(TOTW*w) for w in widths]
    else: ws=[int(TOTW/ncol)]*ncol
    area_h=Inches(4.65) if note else Inches(4.85)
    nrow=len(rows)+1; rh=int(area_h/nrow); y=Inches(1.95)
    x=X0
    for j,htxt in enumerate(headers):
        rect(s,x,y,ws[j],rh,BLUE)
        txt(s,x+Inches(0.14),y,ws[j]-Inches(0.24),rh,[[(htxt,fsize,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)
        x+=ws[j]
    for i,row in enumerate(rows):
        y=int(Inches(1.95)+rh*(i+1)); x=X0
        fill=LIGHT if i%2==0 else WHITE
        for j,cell in enumerate(row):
            rect(s,x,y,ws[j],rh,fill,line=LINE)
            bold=(j==0)
            txt(s,x+Inches(0.14),y,ws[j]-Inches(0.24),rh,[[(cell,fsize-1,INK if not bold else INK,bold)]],anchor=MSO_ANCHOR.MIDDLE)
            x+=ws[j]
    if note:
        txt(s,Inches(0.85),Inches(6.62),Inches(11.7),Inches(0.4),[[(note,13,GREY,False)]],align=PP_ALIGN.CENTER)
    footer(s); return s
def formula_slide(title,panels,kicker=None,accent=BLUE,note=None):
    """Row of formula panels: (heading, formula, caption)."""
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(panels); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.3)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.1); ch=Inches(3.9)
    for i,(hd,formula,cap) in enumerate(panels):
        x=int(X0+(cw+gap)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),col)
        txt(s,x+Inches(0.22),y+Inches(0.28),cw-Inches(0.44),Inches(0.5),[[(hd,16,col,True)]])
        fr=rect(s,x+Inches(0.22),y+Inches(0.95),cw-Inches(0.44),Inches(1.35),NAVY)
        txt(s,x+Inches(0.32),y+Inches(0.95),cw-Inches(0.64),Inches(1.35),
            [[(ln,15,RGBColor(0x9C,0xDC,0xFE),True)] for ln in formula.split("\n")],anchor=MSO_ANCHOR.MIDDLE,space=3)
        txt(s,x+Inches(0.22),y+Inches(2.5),cw-Inches(0.44),ch-Inches(2.6),[[(cap,12.5,GREY,False)]])
    if note:
        txt(s,Inches(0.85),Inches(6.25),Inches(11.7),Inches(0.6),[[(note,14,GREY,False)]],align=PP_ALIGN.CENTER)
    footer(s); return s
def activity_slide(a,topic):
    """ONE compact workflow slide per hands-on activity — scenario line,
    numbered workflow strip, deliverable band, tools+time chips."""
    s=head(slide(),f"Activity {a['num']} — {a['title']}",
           kicker=f"TOPIC {topic['code']} · HANDS-ON ACTIVITY",kcolor=TEAL)
    # badge + duration chip
    rect(s,Inches(10.35),Inches(0.5),Inches(2.13),Inches(0.42),TEAL)
    txt(s,Inches(10.35),Inches(0.5),Inches(2.13),Inches(0.42),
        [[(f"ACTIVITY {a['num']}  ·  {a['minutes']} MIN",13,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    # scenario
    txt(s,Inches(0.85),Inches(1.92),Inches(11.7),Inches(1.15),[[(a["desc"],15,INK,False)]])
    # workflow strip
    steps=a["flow"]; n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.26)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(3.15); ch=Inches(1.95); bd=Inches(0.56)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.08),TEAL)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.24)),bd,bd,TEAL)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.24)),bd,bd,[[(str(i+1),20,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.12),int(y+Inches(0.95)),cw-Inches(0.24),int(ch-Inches(1.05)),[[(st,11.5,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.06)),int(y+ch/2-Inches(0.28)),int(gap+Inches(0.12)),Inches(0.5),
                [[("▶",12,TEAL,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    # deliverable band
    rect(s,Inches(0.85),Inches(5.35),Inches(11.63),Inches(0.92),RGBColor(0xE8,0xF7,0xEE))
    txt(s,Inches(1.1),Inches(5.35),Inches(11.2),Inches(0.92),
        [[("YOU'LL PRODUCE   ",12,RGBColor(0x12,0x7A,0x3E),True),(a["build"],13,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    # tools line
    txt(s,Inches(0.85),Inches(6.38),Inches(11.63),Inches(0.5),
        [[("Tools:  ",12,GREY,True),(a["services"],12,GREY,False),
          (f"     ·     Full step-by-step guide: Learner Guide, Activity {a['num']}",12,GREY,False)]])
    footer(s); return s
def novaspc_slide():
    """Your SPC Lab Tool — NovaSPC: browser-card mock of the ed-tool plus the
    5-step lab workflow every activity follows."""
    s=head(slide(),"Your SPC Lab Tool — NovaSPC",kicker="COURSE ED-TOOL · ALL 12 ACTIVITIES",kcolor=TEAL)
    # browser card
    bx,by,bw,bh=Inches(0.85),Inches(2.0),Inches(6.1),Inches(4.55)
    rect(s,bx,by,bw,bh,WHITE,line=LINE)
    rect(s,bx,by,bw,Inches(0.52),RGBColor(0xEE,0xF2,0xF8))
    for i,c in enumerate([RED,AMBER,TEAL]):
        oval(s,int(bx+Inches(0.18)+Inches(0.28)*i),int(by+Inches(0.17)),Inches(0.18),Inches(0.18),c)
    rect(s,int(bx+Inches(1.2)),int(by+Inches(0.1)),int(bw-Inches(1.5)),Inches(0.34),WHITE,line=LINE)
    txt(s,int(bx+Inches(1.35)),int(by+Inches(0.1)),int(bw-Inches(1.8)),Inches(0.34),
        [[("https://alfredang.github.io/novaspc/",13,TEAL,True)]],anchor=MSO_ANCHOR.MIDDLE)
    # left nav mock
    nx=int(bx+Inches(0.25)); nw=Inches(2.1); ny=int(by+Inches(0.75))
    rect(s,nx,ny,nw,Inches(3.55),LIGHT)
    nav=[("DATA",True),("Data Input",False),("ATTRIBUTE CHARTS",True),("c · u · np · p",False),
         ("VARIABLE CHARTS",True),("X-mR · X̄-R · X̄-s",False),("ANALYSIS",True),("Distribution",False),("Process Capability",False)]
    for i,(t,hd) in enumerate(nav):
        txt(s,nx+Inches(0.15),int(ny+Inches(0.10)+Inches(0.375)*i),nw-Inches(0.3),Inches(0.36),
            [[(t,10 if hd else 11.5,GREY if hd else INK,hd)]])
    # main panel mock — statistics tiles
    mx=int(bx+Inches(2.5)); mw=int(bw-Inches(2.75))
    txt(s,mx,ny,mw,Inches(0.4),[[("X̄-R Chart — Statistics",13,INK,True)]])
    stats=[("10.27","X̄̄ (CL)"),("10.49","X̄ UCL"),("10.05","X̄ LCL"),("0.30","R̄ (CL)"),("0.68","R UCL"),("0","OOC")]
    for i,(v,l) in enumerate(stats):
        r_,c_=i//3,i%3
        tx=int(mx+(mw/3)*c_); ty=int(ny+Inches(0.5)+Inches(0.78)*r_)
        rect(s,tx,ty,int(mw/3-Inches(0.08)),Inches(0.68),LIGHT)
        txt(s,tx,ty+Inches(0.06),int(mw/3-Inches(0.08)),Inches(0.3),[[(v,14,TEAL,True)]],align=PP_ALIGN.CENTER)
        txt(s,tx,ty+Inches(0.36),int(mw/3-Inches(0.08)),Inches(0.26),[[(l,9,GREY,False)]],align=PP_ALIGN.CENTER)
    rect(s,mx,int(ny+Inches(2.2)),mw,Inches(1.3),RGBColor(0xE8,0xF7,0xEE))
    txt(s,mx+Inches(0.15),int(ny+Inches(2.3)),mw-Inches(0.3),Inches(1.1),
        [[("Generate Chart",12,RGBColor(0x12,0x7A,0x3E),True)],
         [("control chart renders here — CL, UCL/LCL and OOC points flagged",10,GREY,False)]],space=3)
    # right: 5-step lab workflow tiles
    steps=[("Get the data","Download the activity's CSV from the Activities folder (labs/data/) — or hand-compute first."),
           ("Upload","Drag the CSV onto Data Input (or load an example dataset) and check the Data Preview."),
           ("Pick the chart","c / u / np / p for attribute data; X-mR / X̄-R / X̄-s for variable data."),
           ("Generate & compare","Click Generate Chart and check the Statistics tiles against YOUR hand computation."),
           ("Analyse & export","Distribution and Process Capability (Cp, Cpk, Pp, Ppk); Export the chart as PNG evidence.")]
    rx=Inches(7.3); rw=Inches(5.2); gy=Inches(0.16); th=int((Inches(4.55)-gy*4)/5)
    for i,(t1,t2) in enumerate(steps):
        y=int(Inches(2.0)+(th+gy)*i); col=PALETTE[i%4]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.09),th,col)
        bd=Inches(0.44)
        oval(s,rx+Inches(0.18),int(y+th/2-bd/2),bd,bd,col)
        txt(s,rx+Inches(0.18),int(y+th/2-bd/2),bd,bd,[[(str(i+1),14,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,rx+Inches(0.78),y,rw-Inches(0.95),th,[[(t1,12.5,col,True)],[(t2,10.5,INK,False)]],anchor=MSO_ANCHOR.MIDDLE,space=2)
    txt(s,Inches(0.85),Inches(6.6),Inches(11.7),Inches(0.4),
        [[("Hand-compute first, then let NovaSPC confirm your numbers — every activity ends with the tool agreeing with you.",13,GREY,False)]],align=PP_ALIGN.CENTER)
    footer(s); return s
def lms_slide():
    """Download Course Material — visual browser-card mock of the LMS portal
    plus numbered how-to tiles (never a bare text link)."""
    s=head(slide(),"Download Course Material",kicker="COURSE PORTAL · LMS/TMS",kcolor=BLUE)
    # browser card
    bx,by,bw,bh=Inches(0.85),Inches(2.0),Inches(6.1),Inches(4.55)
    rect(s,bx,by,bw,bh,WHITE,line=LINE)
    rect(s,bx,by,bw,Inches(0.52),RGBColor(0xEE,0xF2,0xF8))
    for i,c in enumerate([RED,AMBER,TEAL]):
        oval(s,int(bx+Inches(0.18)+Inches(0.28)*i),int(by+Inches(0.17)),Inches(0.18),Inches(0.18),c)
    rect(s,int(bx+Inches(1.2)),int(by+Inches(0.1)),int(bw-Inches(1.5)),Inches(0.34),WHITE,line=LINE)
    txt(s,int(bx+Inches(1.35)),int(by+Inches(0.1)),int(bw-Inches(1.8)),Inches(0.34),
        [[("https://lms-tms.tertiaryinfotech.com",13,BLUE,True)]],anchor=MSO_ANCHOR.MIDDLE)
    # portal body mock
    rect(s,int(bx+Inches(0.25)),int(by+Inches(0.75)),int(bw-Inches(0.5)),Inches(0.75),BLUE)
    txt(s,int(bx+Inches(0.45)),int(by+Inches(0.75)),int(bw-Inches(0.9)),Inches(0.75),
        [[("Tertiary Infotech LMS/TMS",15,WHITE,True)],[("My Courses  ·  Attendance  ·  Assessment",11,RGBColor(0xB3,0xD4,0xFF),False)]],anchor=MSO_ANCHOR.MIDDLE,space=2)
    rows=[("Statistical Process Control (SPC) in Manufacturing",True),
          ("Trainer Slides (PDF)  ⬇",False),("Learner Guide (PDF)  ⬇",False),("Assessment Submission  ↗",False)]
    for i,(rt,hd) in enumerate(rows):
        yy=int(by+Inches(1.68)+Inches(0.68)*i)
        rect(s,int(bx+Inches(0.25)),yy,int(bw-Inches(0.5)),Inches(0.56),LIGHT if not hd else RGBColor(0xE8,0xF0,0xFE))
        txt(s,int(bx+Inches(0.45)),yy,int(bw-Inches(0.9)),Inches(0.56),
            [[(rt,12.5,INK if not hd else BLUE,hd)]],anchor=MSO_ANCHOR.MIDDLE)
    # step tiles right
    steps=[("Sign in","lms-tms.tertiaryinfotech.com — log in with your registered email (OTP or password)."),
           ("Open your course","Select 'Statistical Process Control (SPC) in Manufacturing' under My Courses."),
           ("Download materials","Slides and the Learner Guide — your open-book references for the assessment."),
           ("Submit & survey","Upload your assessment answers and complete the TRAQOM survey QR here.")]
    rx=Inches(7.3); rw=Inches(5.2); gy=Inches(0.2); th=int((Inches(4.55)-gy*3)/4)
    for i,(t1,t2) in enumerate(steps):
        y=int(Inches(2.0)+(th+gy)*i); col=PALETTE[i%4]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.09),th,col)
        bd=Inches(0.5)
        oval(s,rx+Inches(0.2),int(y+th/2-bd/2),bd,bd,col)
        txt(s,rx+Inches(0.2),int(y+th/2-bd/2),bd,bd,[[(str(i+1),16,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,rx+Inches(0.9),y,rw-Inches(1.1),th,[[(t1,14,col,True)],[(t2,11.5,INK,False)]],anchor=MSO_ANCHOR.MIDDLE,space=2)
    txt(s,Inches(0.85),Inches(6.6),Inches(11.7),Inches(0.4),
        [[("All course material is downloaded from the LMS/TMS portal — keep it handy: the final assessment is open book.",13,GREY,False)]],align=PP_ALIGN.CENTER)
    footer(s); return s

# ============================================================ BUILD
cover()

# ---------------- ADMIN ----------------
mark("admin")
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="TRAQOM · SSG DIGITAL ATTENDANCE")
trainer_slide("YOUR TRAINER · GENERAL","Your Trainer","General Trainer template —\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",C.TRAINER,"Principal Trainer\nTertiary Infotech Academy Pte Ltd",
 [("Role","Principal Trainer, Tertiary Infotech Academy Pte Ltd"),
  ("Background","PhD — 20+ years of industry and training experience in manufacturing quality, statistics and data analytics."),
  ("Delivers","WSQ courses on statistical process control, Lean Six Sigma, data science and AI."),
  ("Founder","Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
 initials="AA",accent=BLUE)
content("Let's Know Each Other",[
 "Your name and organisation / role.",
 "Your experience with SPC, quality control or manufacturing processes (if any).",
 "One process at your workplace you would like to bring under statistical control."],kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Participate actively — no question is too small.",
 "Mutual respect: agree to disagree.","One conversation at a time.",
 "Be punctual; return from breaks on time.","75% attendance is required."],
 kicker="HOUSEKEEPING",cols=2,size=15)
lms_slide()
novaspc_slide()
two_col("Lesson Plan — 1 Day, 9:00am–6:00pm",[
 ("Morning (AM attendance)",0,True),
 ("Welcome, introductions, learning outcomes",1),
 ("Topic 1: Introduction to SPC + Activities 1–2",1),
 ("Tea break (15 min)",1),
 ("Topic 2: Control Charts + Activities 3–9",1),
 ("Lunch break 1:00–2:00pm",1)],
 [("Afternoon (PM attendance)",0,True),
 ("Topic 3: Setup SPC + Activity 10",1),
 ("Topic 4: Process Control Capabilities + Activity 11",1),
 ("Tea break (15 min)",1),
 ("Topic 5: OOC & Follow-Up Actions + Activity 12",1),
 ("TRAQOM survey · Briefing · Assessment attendance",1),
 ("WA(Q&A) 80 min + Oral Questioning 10 min",1)],
 kicker="SCHEDULE",lhead="Morning — foundations & charts",rhead="Afternoon — setup, capability, OOC",
 note="8 training hours: 6.5 h classroom facilitation + 1.5 h assessment · 1-hour lunch · tea breaks counted within training time.")
table_slide("Skills Framework — Quality Process Control",
 ["TSC element","Detail"],
 [("TSC Title / Code",f"{C.TSC_TITLE}  ·  {C.TSC_CODE}")]+
 [(f"Ability {c}",d) for c,d in C.TSC_ABILITIES]+
 [(f"Knowledge {c}",d) for c,d in C.TSC_KNOWLEDGE],
 kicker="WSQ ALIGNMENT",widths=[0.28,0.72],
 note="Every topic, activity and assessment item in this course maps to these TSC abilities and knowledge factors.")
tile_grid("Learning Outcomes",
 [(lo.split(": ")[0],lo.split(": ")[1]) for lo in C.LEARNING_OUTCOMES],
 kicker="BY THE END OF THIS COURSE",cols=1,size=14)
two_col("Course Outline — Five Topics",[
 ("Topic 1 — Introduction to SPC",0,True),
 ("Variation · SPC & why manufacturing uses it · data types · distributions · CLT",1),
 ("Topic 2 — Control Charts",0,True),
 ("Subgroups · Xbar-R · Xbar-S · I-MR · p · np · c · u",1),
 ("Topic 3 — Setup SPC",0,True),
 ("Construct charts · control limits · interpretation",1)],
 [("Topic 4 — Process Control Capabilities",0,True),
 ("Capability vs control · Cp & Cpk · GR&R",1),
 ("Topic 5 — OOC & Follow-Up Actions",0,True),
 ("SPC rules · out of control · root cause analysis · follow-up",1),
 ("12 hands-on activities",0,True),
 ("Every topic is practised on real numbers you compute yourself",1)],
 kicker="THE JOURNEY",lhead="Morning",rhead="Afternoon")
mark("briefing")
content("Briefing for Assessment",[
 "Place phones and other materials under the table or on the floor.",
 "No photos or recording of assessment scripts.","No discussion during the assessment.",
 "Use a black/blue pen for hard-copy assessments.","No liquid paper / correction tape.",
 "Assessment scripts are collected when the time is up."],kicker="BEFORE THE ASSESSMENT")
content("Final Assessment",[
 C.ASSESSMENT["written"], C.ASSESSMENT["practical"],
 "Format: Open Book — slides, Learner Guide and approved materials only.",
 C.ASSESSMENT["note"],
 "An appeal process is available if you disagree with the assessment outcome."],kicker="FINAL ASSESSMENT")
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA(Q&A) 80 min, then OQ 10 min — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")

# ---------------- TOPIC 1 ----------------
T=C.TOPICS[0]; mark("topic1")
section(f"TOPIC {T['code']}",T["title"],T["code"],T["subtitle"])
tile_grid("Sources of Variation — the 6 Ms",[
 ("Material","Quality and consistency of incoming lots and suppliers."),
 ("Machine","Settings, wear and condition of the equipment."),
 ("Manpower (Man)","Operator standards, skill and consistency."),
 ("Method","How the processing is actually carried out."),
 ("Measurement","Gauge repeatability and reproducibility."),
 ("Environment","Temperature, humidity and surroundings.")],
 kicker="NO TWO THINGS ARE IDENTICAL",cols=2,size=15)
tile_grid("Types of Variation",[
 ("Within unit / lot","Positional variation inside one unit or lot."),
 ("Unit to unit","Differences between consecutive units."),
 ("Lot to lot","Differences between production lots."),
 ("Line to line","Differences between production lines."),
 ("Across time","Drift and change from hour to hour, day to day."),
 ("Measurement","Gauge repeatability & reproducibility (GR&R).")],
 kicker="WHERE THE VARIATION LIVES",cols=2,size=15)
two_col("Random vs Non-Random Variation",[
 ("Common (random) causes",0,True),
 ("Inherent to the process — non-assignable",1),
 ("Many small sources acting together",1),
 ("Only management can reduce them: redesign the process, new methods, new equipment",1),
 ("A process with ONLY common causes is IN statistical control",1)],
 [("Special (assignable) causes",0,True),
 ("Unforeseen, non-periodic events",1),
 ("A specific, findable reason: a bad lot, a worn tool, a skipped setup",1),
 ("The goal of SPC: detect them, find them, remove them",1),
 ("A process WITH special causes is OUT of control",1)],
 kicker="TWO FAMILIES OF CAUSES",lhead="Common causes — the system",rhead="Special causes — events",
 note="Treatment differs: special causes are investigated and removed at the line; common causes need management action on the system itself.")
tile_grid("What is Statistical Process Control?",[
 ("Cost-effective quality","SPC monitors and forecasts, predicting problems BEFORE they occur."),
 ("Reactive → preventive","The focus shifts from inspecting defects out to preventing them."),
 ("Assess variability","The goal of SPC is to assess and reduce variability in a process."),
 ("'In control' defined","Products vary consistently within expected limits over time — a stable bell curve.")],
 kicker="SPC IN ONE SLIDE",cols=2,size=15)
img_full("Manufacturing With vs Without SPC","spc_vs_no_spc.png",kicker="WHY USE SPC",
 caption="Prevention beats inspect-and-sort: the control chart signals BEFORE defects are produced, so correction is cheap and the product ships right the first time.")
two_col("Two Types of Manufacturing Data",[
 ("Continuous / variable data",0,True),
 ("Measured on a scale: thickness, height, length, weight, viscosity",1),
 ("Modelled by the NORMAL distribution",1),
 ("Charts: Xbar-R, Xbar-S, I-MR",1)],
 [("Discrete / attribute data",0,True),
 ("Counted categories: pass/fail, defect counts",1),
 ("Defectives → BINOMIAL → p, np charts",1),
 ("Defects → POISSON → c, u charts",1)],
 kicker="DATA DECIDES THE CHART",lhead="Variable — you measure it",rhead="Attribute — you count it",
 note="This split drives everything in Topic 2: the data type and the subgroup size together select the control chart.")
img_points("Binomial Distribution — Counting Defectives","dist_binomial.png",[
 ("Pass or fail","Each inspected unit has exactly two outcomes — defective or not."),
 ("Mean and variance","Mean = np; variance = np(1 − p) for sample size n, proportion p."),
 ("Feeds p and np charts","Control limits for defectives charts come from this distribution.")],
 kicker="DISTRIBUTIONS · 1 OF 3",img_w=6.9)
img_points("Poisson Distribution — Counting Defects","dist_poisson.png",[
 ("Defects per unit","Counts of events in an interval — several defects can share one unit."),
 ("One parameter λ","λ is the expected defects per unit; mean = variance = λ."),
 ("Feeds c and u charts","Control limits for defects charts use σ = √λ.")],
 kicker="DISTRIBUTIONS · 2 OF 3",img_w=6.9)
img_points("Normal Distribution — the Bell Curve","dist_normal.png",[
 ("Measurement data","Continuous quality characteristics cluster symmetrically around the mean."),
 ("68 – 95 – 99.7","±1σ holds 68.3%, ±2σ holds 95.4%, ±3σ holds 99.73% of a stable process."),
 ("Only 0.27% beyond ±3σ","So a point outside ±3σ almost certainly signals a special cause — Shewhart's key insight.")],
 kicker="DISTRIBUTIONS · 3 OF 3",img_w=6.9)
img_points("Central Limit Theorem","clt.png",[
 ("Means become normal","Sample means are normally distributed for large n — whatever the population looks like."),
 ("n ≈ 30 is enough","Around 30 observations assure the bell shape; subgroup averages get there much faster."),
 ("Why charts work","Xbar charts plot AVERAGES — the CLT licenses ±3σ limits even for non-normal processes.")],
 kicker="THE THEOREM BEHIND THE CHARTS",img_w=7.4)
formula_slide("Sample Statistics — Location and Dispersion",[
 ("Mean (location)","X̄ = ΣXᵢ / n","The central tendency of the sample — the 'where' of the process. Grand average X̿ averages the subgroup means."),
 ("Range (dispersion)","R = Xmax − Xmin","The simplest spread measure — cheap to compute on the shop floor; used by the R chart for subgroups of 10 or less."),
 ("Std deviation (dispersion)","s = √( Σ(Xᵢ−X̄)² / (n−1) )","The efficient spread measure — used by the S chart when subgroups exceed 10, and by capability indices.")],
 kicker="THE STATISTICS EVERY CHART PLOTS",
 note="Central tendency: mean.  Dispersion: range and standard deviation.  Every control chart pairs one location statistic with one dispersion statistic.")
mark("act1"); activity_slide(ACTIVITIES[0],T)
mark("act2"); activity_slide(ACTIVITIES[1],T)

# ---------------- TOPIC 2 ----------------
T=C.TOPICS[1]; mark("topic2")
section(f"TOPIC {T['code']}",T["title"],T["code"],T["subtitle"])
tile_grid("Controlling a Process",[
 ("Get to normal first","Remove assignable causes — you don't want to 'control' a broken process."),
 ("Monitor at intervals","Sample often enough to catch deterioration before rejects are made."),
 ("Use control charts","Plot the property in time order against statistically derived limits."),
 ("Act before defects","Take corrective action as soon as the chart signals — that is the whole point.")],
 kicker="WHAT 'IN CONTROL' TAKES",cols=2,size=15)
img_points("Origin of Control Charts — Dr Walter A. Shewhart","control_chart_anatomy.png",[
 ("Bell Labs, 1924","Shewhart invented the Xbar-R chart — the first control chart — for manufacturing."),
 ("The ±3σ insight","For a normal process only 0.27% of points fall beyond ±3σ — rare enough to signal a real problem."),
 ("Rotate and plot","Rotate the normal curve 90°, draw limits at ±3σ, plot subgroups in time order — a control chart."),
 ("Three elements","Time-series graph + centre line (location) + UCL/LCL equidistant limits (dispersion).")],
 kicker="THE FATHER OF CONTROL CHARTS",img_w=7.2)
img_points("Why 3 Sigma? — the Economic Argument","three_sigma_cost.png",[
 ("Tighter limits (1–2σ)","Many false alarms — you chase ghosts and pay α costs investigating a stable process."),
 ("Wider limits (4–5σ)","Real shifts slip through — β costs mount as defects reach the customer."),
 ("±3σ minimises total cost","Shewhart's choice balances both — the most economical control limit.")],
 kicker="NOT ARBITRARY — ECONOMIC",img_w=7.2)
tile_grid("Rational Subgrouping",[
 ("Homogeneous within","Choose subgroups so differences WITHIN a subgroup are only common cause."),
 ("Opportunity between","Give maximum opportunity for variation FROM ONE SUBGROUP TO ANOTHER."),
 ("Small and frequent","Typical initial study: 4–5 units every hour — signals a shift before much defective product is made."),
 ("Enough history","Collect ≥ 100 observations (e.g. 25 subgroups of 4) before firming up limits."),
 ("Not too large","A subgroup spanning a whole day averages out the very changes you want to see."),
 ("Relax as it stabilises","As the process demonstrates stability, reduce subgroup size and frequency.")],
 kicker="THE KEY IDEA IN SHEWHART CHARTS",cols=2,size=14)
table_slide("Choosing a Control Chart",
 ["Data type","Statistic charted","Subgroup size","Chart"],
 [("Continuous (variable)","Mean + range","2 – 10","Xbar-R"),
  ("Continuous (variable)","Mean + std deviation","> 10","Xbar-S"),
  ("Continuous (variable)","Individuals + moving range","1","I-MR (X-MR)"),
  ("Attribute — defectives","Proportion defective","Varies","p chart"),
  ("Attribute — defectives","Number defective","Constant","np chart"),
  ("Attribute — defects","Defects per unit","Varies","u chart"),
  ("Attribute — defects","Defect count","Constant","c chart")],
 kicker="DATA TYPE + SUBGROUP SIZE DECIDE",widths=[0.28,0.28,0.18,0.26],
 note="Variable data uses PAIRED charts (location on top, dispersion below); attribute data uses a single chart.")
img_points("Xbar-R Chart — Mean and Range","xbar_r_chart.png",[
 ("When to use","Continuous data in time order, normally distributed, subgroup size 2–10."),
 ("Xbar panel","Subgroup MEANS — monitors process location between subgroups."),
 ("R panel","Subgroup RANGES — monitors within-subgroup spread over time."),
 ("Read R first","Xbar limits are derived from R-bar — an out-of-control R chart invalidates the Xbar limits.")],
 kicker="VARIABLE CHARTS · THE WORKHORSE",img_w=6.6)
flow_h("How to Interpret an Xbar-R Chart",[
 "Examine the R chart FIRST — its stability underwrites the Xbar limits",
 "R out of control? Stop, find the special cause, remove those subgroups",
 "R in control → now read the Xbar chart against its limits",
 "Judge against CONTROL limits — never against customer spec limits",
 "Both panels stable? Only now may capability studies begin"],
 kicker="ALWAYS IN THIS ORDER",color=VIOLET,
 note="Specification limits are the customer's voice; control limits are the process's voice — they never belong on the same panel.")
mark("act3"); activity_slide(ACTIVITIES[2],T)
img_points("Xbar-S Chart — Mean and Standard Deviation","xbar_s_chart.png",[
 ("When to use","Same conditions as Xbar-R, but subgroup size ABOVE 10."),
 ("Why S beats R there","With large subgroups the range wastes information; S uses every reading."),
 ("Same discipline","Read the S chart first — Xbar limits derive from S-bar."),
 ("Same data, new lens","Our five thickness subgroups give S values 0.129–0.183 with S-bar 0.135.")],
 kicker="VARIABLE CHARTS · LARGE SUBGROUPS",img_w=6.6)
mark("act4"); activity_slide(ACTIVITIES[3],T)
img_points("I-MR Chart — Individuals and Moving Range","imr_chart.png",[
 ("Subgroup of ONE","Slow production, expensive or destructive tests, batch processes."),
 ("I panel","Every individual reading against limits from the moving range."),
 ("MR panel","|current − previous| — short-term variation between consecutive readings."),
 ("Batch example","Eight viscosity batches: I-bar 25.39 cP, MR-bar 0.33 cP.")],
 kicker="VARIABLE CHARTS · N = 1",img_w=6.6)
table_slide("Xbar-R vs Xbar-S vs I-MR",
 ["Chart","Statistics","Subgroup size","Variation monitored"],
 [("Xbar","Mean","2+","Between subgroups (location)"),
  ("R","Range","2 – 10","Within subgroup (spread)"),
  ("S","Std deviation","> 10","Within subgroup (spread)"),
  ("I","Individuals","1","Process level over time"),
  ("MR","Moving range","1","Reading-to-reading spread")],
 kicker="THE VARIABLE CHART FAMILY",widths=[0.16,0.24,0.22,0.38])
mark("act5"); activity_slide(ACTIVITIES[4],T)
img_full("The Four Attribute Charts — One Family Portrait","attribute_charts.png",kicker="ATTRIBUTE CHARTS",
 caption="Defectives (binomial): p plots the proportion, np the count with constant n.  Defects (Poisson): u plots defects per unit, c the count with a constant inspection unit.")
two_col("p and np Charts — Defectives",[
 ("p chart — proportion defective",0,True),
 ("p = defectives ÷ sample size, per batch",1),
 ("Sample size MAY vary — limits adjust per point",1),
 ("Batch example: p-bar = 10/125 = 0.08, UCL ≈ 0.24",1),
 ("Uses: sudden change detection, before/after comparison, stability checks",1)],
 [("np chart — number defective",0,True),
 ("Plots the COUNT of defective units",1),
 ("Requires a CONSTANT sample size",1),
 ("Daily example: np-bar = 4.0, UCL = 4 + 3·√(np̄(1−p̄)) ≈ 9.75",1),
 ("Same binomial engine as p — just unscaled",1)],
 kicker="BINOMIAL CHARTS",lhead="p — proportions",rhead="np — counts",
 note="An 'item' can be anything you chart: gadgets from a line, wait times, delivery times — anything with a pass/fail outcome.")
mark("act6"); activity_slide(ACTIVITIES[5],T)
mark("act7"); activity_slide(ACTIVITIES[6],T)
two_col("c and u Charts — Defects",[
 ("c chart — defect count",0,True),
 ("Counts defects where one item can carry many",1),
 ("Requires a CONSTANT inspection unit",1),
 ("TV example: c-bar = 6.5 pixels, UCL = c̄ + 3·√c̄ ≈ 14.15, LCL floored at 0",1)],
 [("u chart — defects per unit",0,True),
 ("u = defects ÷ units inspected, per sample",1),
 ("Inspection quantity MAY vary — limits step with each n",1),
 ("PCB example: u-bar = 80/120 ≈ 0.67 defects per board",1)],
 kicker="POISSON CHARTS",lhead="c — constant unit",rhead="u — varying units",
 note="Both monitor process stability over time so instabilities can be identified and corrected — the Poisson σ = √mean does the work.")
mark("act8"); activity_slide(ACTIVITIES[7],T)
mark("act9"); activity_slide(ACTIVITIES[8],T)

# ---------------- TOPIC 3 ----------------
T=C.TOPICS[2]; mark("topic3")
section(f"TOPIC {T['code']}",T["title"],T["code"],T["subtitle"])
flow_h("SPC Implementation Pipeline",[
 "Select the critical characteristic and the right chart type",
 "Plan rational subgroups — size, frequency, duration",
 "Collect 20+ subgroups of data in time order",
 "Compute trial centre lines and ±3σ control limits",
 "Plot, interpret, remove assignable causes — then run and improve"],
 kicker="FROM CHARACTERISTIC TO RUNNING CHART",
 note="Trial limits are provisional: recompute after removing assignable-cause subgroups, and review as the process changes.")
formula_slide("How Control Limits Are Computed",[
 ("The principle","UCL = X̿ + 3σ/√n\nCL  = X̿\nLCL = X̿ − 3σ/√n","Estimate σ from the data, multiply by 3, scale by √n for averages — the CLT at work."),
 ("Xbar-R shortcut","UCL(X̄) = X̿ + A2·R̄\nLCL(X̄) = X̿ − A2·R̄","Chart constants fold 3/(d2·√n) into one number A2 you look up by subgroup size."),
 ("Range chart","UCL(R) = D4·R̄\nLCL(R) = D3·R̄","D3 = 0 for n ≤ 6 — small-subgroup R charts have no lower limit.")],
 kicker="±3σ, OPERATIONALISED",
 note="Attribute charts use their distribution's σ directly:  p̄ ± 3√(p̄(1−p̄)/n)  ·  np̄ ± 3√(np̄(1−p̄))  ·  c̄ ± 3√c̄  ·  ū ± 3√(ū/n).")
table_slide("Control Chart Constants",
 ["Subgroup size n","A2 (Xbar)","D3 (R lower)","D4 (R upper)","A3 (Xbar-S)"],
 [("2","1.880","—","3.267","2.659"),
  ("3","1.023","—","2.574","1.954"),
  ("4","0.729","—","2.282","1.628"),
  ("5","0.577","—","2.114","1.427"),
  ("6","0.483","—","2.004","1.287"),
  ("7","0.419","0.076","1.924","1.182"),
  ("10","0.308","0.223","1.777","0.975")],
 kicker="LOOK UP BY SUBGROUP SIZE",widths=[0.24,0.19,0.19,0.19,0.19],
 note="n = 4 (our labs): A2 = 0.729, D3 = 0, D4 = 2.282.  No lower range limit exists for subgroups of 6 or less.")
tile_grid("Constructing an Xbar-R Chart",[
 ("1 · Record","Collect subgroup observations in time order — 20+ subgroups before limits firm up."),
 ("2 · Compute","Xbar and R for every subgroup; then X̿ (grand average) and R̄."),
 ("3 · R limits","UCL(R) = D4·R̄, LCL(R) = D3·R̄ — plot the R chart and confirm it is in control."),
 ("4 · Xbar limits","UCL/LCL(X̄) = X̿ ± A2·R̄ — plot the Xbar chart with its centre line."),
 ("5 · Interpret","Points against CONTROL limits, R chart first; investigate any signal."),
 ("6 · Maintain","Plot new data live; do NOT recalculate limits unless the process permanently changes.")],
 kicker="SIX MOVES, ONE CHART",cols=2,size=13)
mark("act10"); activity_slide(ACTIVITIES[9],T)
two_col("Revising and Updating Control Limits",[
 ("When to REVISE",0,True),
 ("A sustained level change on either chart — usually 20+ points",1),
 ("Regular review periods: weekly, monthly, or every 50 samples",1),
 ("Some users centre Xbar on a target (nominal spec) instead of X̿",1)],
 [("Triggers to UPDATE",0,True),
 ("Change of supplier for a critical material",1),
 ("Change in process machinery or engineering change orders",1),
 ("Introduction of new operators",1),
 ("Change in sample size",1)],
 kicker="LIMITS MUST REFLECT TODAY'S PROCESS",lhead="Revision policy",rhead="Update triggers",
 note="The chart must always reflect the PRESENT condition of the process — stale limits alarm on ghosts or sleep through real shifts.")
tile_grid("Interpreting Control Charts — the Ground Rules",[
 ("Two panels for variables","Averages chart = between-subgroup variation; R/S chart = within-subgroup variation."),
 ("Centre line vs objective","First check the process mean is where the business needs it — capable is a separate question."),
 ("Inside ≠ ignore","Trends, shifts and patterns INSIDE the limits are special causes too."),
 ("Control ≠ conformance","'In control' means consistent — not that the product meets spec."),
 ("Remove and redo","Points traced to a cause are removed and the limits recalculated."),
 ("Charting alone changes nothing","A signal without investigation and action is just decoration — form a team and act.")],
 kicker="READ THE CHART LIKE A PRACTITIONER",cols=2,size=13)

# ---------------- TOPIC 4 ----------------
T=C.TOPICS[3]; mark("topic4")
section(f"TOPIC {T['code']}",T["title"],T["code"],T["subtitle"])
two_col("Capability vs Control — Different Questions",[
 ("CONTROL — is it stable?",0,True),
 ("All special causes removed; only common-cause variation remains",1),
 ("Judged on the CONTROL CHART: no points out, no trends, no patterns",1),
 ("The process is predictable — it will keep doing what it is doing",1)],
 [("CAPABILITY — does it fit spec?",0,True),
 ("Nearly 100% of output inside the customer's specification limits",1),
 ("Judged with Cp / Cpk against LSL and USL",1),
 ("A process can be perfectly in control and still NOT capable — then improve or redesign it",1)],
 kicker="STABLE FIRST, CAPABLE SECOND",lhead="Voice of the process",rhead="Voice of the customer",
 note="Capability studies are meaningless on an unstable process — establish statistical control before computing Cp or Cpk.")
img_full("Capability = Accuracy AND Precision","accuracy_precision.png",kicker="THE TARGET VIEW",
 caption="Precision is a small spread (low σ); accuracy is being on target (mean at spec centre). Cp sees only precision — Cpk sees both.")
img_full("Capability = How Few Defects Escape the Spec","capability_defects.png",kicker="THE DISTRIBUTION VIEW",
 caption="Same spec limits, two processes: when the ±3σ spread fits inside LSL–USL almost nothing is out of spec; when it spills, the red tails are your defect rate.")
img_points("Process Capability Ratios — Cp and Cpk","cp_cpk.png",[
 ("Cp — precision only","Cp = (USL − LSL) / 6σ.  Cp > 1: potentially capable; Cp < 1: development needed."),
 ("Cpk — precision + accuracy","Cpk = min(USL − μ, μ − LSL) / 3σ — punishes an off-centre mean."),
 ("Read the gap","Cp 1.67 vs Cpk 1.00 = pure mis-centring; re-centre the mean and Cpk rises to Cp."),
 ("Common targets","Cpk ≥ 1.00 capable · ≥ 1.33 standard requirement · ≥ 1.67 for critical characteristics.")],
 kicker="THE TWO NUMBERS EVERY CUSTOMER ASKS FOR",img_w=6.9)
two_col("Control Limits vs Specification Limits",[
 ("Specification limits (LSL, USL)",0,True),
 ("Set by DESIGN / customer requirements",1),
 ("Apply to INDIVIDUAL product values",1),
 ("External to the process — they don't move when the process does",1)],
 [("Control limits (LCL, UCL)",0,True),
 ("Derived from the PROCESS's own variability",1),
 ("Apply to SAMPLE statistics — subgroup means and ranges",1),
 ("Recomputed only when the process permanently changes",1)],
 kicker="NEVER CONFUSE THE TWO",lhead="Customer's voice",rhead="Process's voice",
 note="Drawing spec limits on an Xbar chart is a classic error — averages spread less than individuals (σ/√n), so the comparison misleads.")
mark("act11"); activity_slide(ACTIVITIES[10],T)
tile_grid("Gauge Repeatability & Reproducibility (GR&R)",[
 ("Why it matters","All measurement data contains error; SPC needs accurate, precise data to act on."),
 ("What GR&R does","Quantifies how much observed variation comes from the MEASUREMENT SYSTEM vs the process."),
 ("Repeatability","Equipment variation — same operator, same gauge, same part, repeated."),
 ("Reproducibility","Operator variation — different operators, same gauge, same part."),
 ("What it reveals","Measurement vs process variation share · operator influence · ability to discriminate parts."),
 ("Act on it","A gauge that can't tell good parts from bad must be fixed before charting its data.")],
 kicker="TRUST THE GAUGE BEFORE THE CHART",cols=2,size=13)
tile_grid("Gage Capability — Five Performance Factors",[
 ("Calibration","Evaluate and adjust against a traceable standard source."),
 ("Resolution","The smallest increment the gauge can distinguish — fine enough for the tolerance."),
 ("Stability","Variation in measurements of the same part over TIME."),
 ("Bias","Systematic offset between observed values and the standard value."),
 ("Linearity","Change in bias across the measuring range — accurate at 1 mm but not at 10 mm.")],
 kicker="WHERE GAUGE ERROR COMES FROM",cols=2,size=14)

# ---------------- TOPIC 5 ----------------
T=C.TOPICS[4]; mark("topic5")
section(f"TOPIC {T['code']}",T["title"],T["code"],T["subtitle"])
two_col("Objectives of Control Charts & Signals of Trouble",[
 ("The one primary purpose",0,True),
 ("Detect assignable causes causing process shifts",1),
 ("…so corrective action removes them BEFORE many non-conforming units are made",1),
 ("Secondary: reduce variability; estimate process parameters and capability",1)],
 [("Indicators of instability",0,True),
 ("PRIMARY: any point outside a control limit",1),
 ("SECONDARY: non-random patterns inside the limits",1),
 ("Shift · trend · stratification · mixture · periodicity",1)],
 kicker="WHAT THE CHART IS FOR",lhead="Objectives",rhead="Indicators",
 note="When a signal appears: mark it on the chart, investigate, and document what you found, the cause, and how it was corrected.")
tile_grid("Common Causes of Instability",[
 ("People","New workers, skills or motivation changes, operator rotation and fatigue."),
 ("Methods & standards","New methods, changed inspection methods or standards, wrong limit calculations."),
 ("Materials & machines","New raw-material lots, new machines, fluctuating machine settings."),
 ("Environment & wear","Systematic temperature changes, maintenance schedule, tool wear.")],
 kicker="WHAT USUALLY BROKE",cols=2,size=15)
img_points("Control Chart Zones","zones.png",[
 ("Three 1σ bands per side","Zone C (0–1σ), Zone B (1σ–2σ), Zone A (2σ–3σ) — mirrored below the CL."),
 ("Why zones exist","They turn 'looks odd' into TESTABLE rules about where points should rarely fall."),
 ("Random looks boring","A stable process hugs the CL: most points in C, few in B, almost none in A.")],
 kicker="THE GRID BEHIND THE RULES",img_w=7.2)
img_full("SPC Rules — the Western Electric Tests","spc_rules.png",kicker="NON-RANDOM PATTERNS",
 caption="Rule 1: one point beyond ±3σ.  Rule 2: 2 of 3 consecutive in Zone A (same side).  Rule 3: run of 8 on one side of the CL — a SHIFT.  Rule 4: 6+ steadily rising or falling — a TREND.")
table_slide("Rule Violations → Likely Causes",
 ["Signal on the chart","What it usually means","Typical causes"],
 [("Point beyond ±3σ","A sudden special cause","Broken tool, wrong material lot, measurement blunder"),
  ("Run of 8 one side (shift)","The process level moved","New operator/method/material, machine reset"),
  ("Trend of 6+ (drift)","Gradual change","Tool wear, temperature drift, operator fatigue"),
  ("Stratification (hugging CL)","Too-good-to-be-true spread","Mixed lots averaged out, wrong subgrouping, limits mis-calculated"),
  ("Mixture / periodicity","Two processes or a cycle","Alternating machines/shifts, systematic environment cycle")],
 kicker="FROM PATTERN TO SUSPECT LIST",widths=[0.28,0.30,0.42],fsize=12,
 note="Annotate the assignable cause on the chart and EXCLUDE those points from capability calculations.")
flow_h("Out-of-Control Procedure",[
 "Detect the signal — a rule fires on the chart",
 "Mark and annotate it; hold or contain suspect product",
 "Investigate and identify the root cause",
 "Correct the cause and verify on the chart",
 "Document everything; resume charting (recalculate limits only for permanent change)"],
 kicker="WHEN A RULE FIRES",color=RED,
 note="Document how you investigated, what you learned, the cause, and how it was corrected — the annotation is part of the record.")
img_points("Pareto Analysis — Find the Vital Few","pareto.png",[
 ("Rank the causes","Bar-chart defect causes from biggest to smallest contributor."),
 ("The 80/20 pattern","A vital few causes typically produce most of the defects."),
 ("Work top-down","Attack the tallest bars first — maximum defect reduction per fix.")],
 kicker="ROOT CAUSE ANALYSIS · TOOL 1",img_w=7.0)
img_points("Ishikawa Fishbone — Organise the Causes","fishbone.png",[
 ("Effect at the head","State the problem precisely — 'coating thickness OOC, 3 points Zone A'."),
 ("6 M bones","Man, Machine, Material, Method, Measurement, Environment."),
 ("Ask why, repeatedly","Drill each promising bone down to an actionable root cause.")],
 kicker="ROOT CAUSE ANALYSIS · TOOL 2",img_w=7.4)
mark("act12"); activity_slide(ACTIVITIES[11],T)
img_full("Improve the Capability — PDSA and DMAIC","pdsa_dmaic.png",kicker="FOLLOW-UP FRAMEWORKS",
 caption="After the OOC is contained, drive systematic improvement: Deming's Plan-Do-Study-Act cycle, or Six Sigma's DMAIC — where Control means holding the gains WITH SPC.")

# ---------------- CLOSE ----------------
mark("close")
section("WRAP-UP","Course Summary & Next Steps","")
tile_grid("What You Achieved",[
 ("Statistical foundations","Variation, distributions, CLT, sample statistics — the numbers behind the charts. (LO1)"),
 ("Chart selection","Matched data type + subgroup size to Xbar-R/S, I-MR, p, np, c, u. (LO2)"),
 ("SPC setup","Built control charts with correct centre lines and ±3σ limits from chart constants. (LO3)"),
 ("Capability verified","Computed and interpreted Cp and Cpk; understood GR&R and gauge performance. (LO4)"),
 ("OOC response","Read SPC rules, ran root cause analysis, and planned follow-up actions. (LO5)")],
 kicker="LEARNING OUTCOMES DELIVERED",cols=1,size=13)
tile_grid("Recommended Follow-On Courses",
 [(t.replace("WSQ - ",""),"WSQ funded · Tertiary Infotech Academy") for t in C.RECOMMENDED],
 kicker="CONTINUE THE QUALITY JOURNEY",cols=1,size=13)
content("Support",[
 "If you have any enquiries during and after the class, contact us:",
 "Email: enquiry@tertiaryinfotech.com",
 "Tel: +65 6100 0613",
 "Website: www.tertiarycourses.com.sg"],kicker="WE'RE HERE TO HELP")
mark("assessment")
content("Assessment",[
 "Written Assessment WA(Q&A) — 80 minutes.  Oral Questioning (OQ) — 10 minutes.",
 "Open book: slides, Learner Guide and approved materials only.",
 "Remember to take the Assessment digital attendance (TRAQOM).",
 "Submit your completed answers on the LMS at https://lms-tms.tertiaryinfotech.com/."],kicker="FINAL ASSESSMENT")
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA(Q&A) 80 min, then OQ 10 min — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="TRAQOM · SSG DIGITAL ATTENDANCE")
big_statement("Thank You!","You can now set up SPC, read its charts, and act on what they say — keep your processes in control.","SEE YOU ON THE LINE",color=TEAL)
mark("end")

OUT=os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
prs.save(OUT)
with open(os.path.join(HERE,"slide_map.json"),"w") as f: json.dump(SLIDE_MAP,f,indent=1)
print(f"Saved {OUT}  ({PAGE['n']} slides)")
print("Slide map:",json.dumps(SLIDE_MAP))
