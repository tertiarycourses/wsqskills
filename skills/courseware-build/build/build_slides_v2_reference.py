#!/usr/bin/env python3
"""Generate the Gemini Agent ADK course slide deck (all-white Tertiary house style).

Design helpers are the same set used by the tertiary-course-slides skill that
produced the n8n reference deck (cover, section, content, two_col, cards3,
big_statement, step_slide, test_slide, brk). Content is driven entirely by
course_data.py + data_domainN.py so the deck stays 100% aligned with the LP,
LG and labs.
"""
import os, sys, copy, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import ChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
from data_domain1 import DOMAIN1
from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3
from data_domain4 import DOMAIN4
ACTIVITIES = DOMAIN1 + DOMAIN2 + DOMAIN3 + DOMAIN4

def _find_repo(start):
    """Locate the course repo (a dir containing both courseware/ and labs/).
    Env COURSE_REPO overrides. Keeps the build working wherever the skill lives."""
    env = os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env):
        return env
    d = start
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "labs")):
            return d
    return os.path.dirname(os.path.dirname(HERE))
REPO = _find_repo(HERE)
ASSETS = os.path.join(os.path.dirname(HERE), "assets")   # co-located with the skill

# ---------------- palette (matches reference) ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=color if lvl==0 else GREY
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

# ---------------- motion: restrained transitions + build animation ----------------
# House rule: ONE transition family for the whole deck (morph is unavailable in the
# OOXML PowerPoint 2010 transition set that python-pptx can emit, so we use a short
# push/fade pair). Content slides fade; section dividers push. Nothing else moves.
P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"

def _transition(s, kind="fade", speed="med"):
    """Attach a restrained slide transition. kind: fade | push | wipe."""
    sld = s._element
    for old in sld.findall(qn("p:transition")):
        sld.remove(old)
    tr = etree.SubElement(sld, qn("p:transition"))
    tr.set("spd", speed)
    tr.set("advClick", "1")
    if kind == "fade":
        etree.SubElement(tr, qn("p:fade"))
    elif kind == "push":
        el = etree.SubElement(tr, qn("p:push")); el.set("dir", "l")
    elif kind == "wipe":
        el = etree.SubElement(tr, qn("p:wipe")); el.set("dir", "r")
    # keep the transition last in the slide element (schema order)
    sld.append(tr)
    return tr

# Appear-on-click build for a list of shape ids — used ONLY on process maps so the
# trainer can reveal one stage at a time. No spins, no flying, no sound.
_TIMING = """<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
<p:childTnLst><p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq">
<p:childTnLst>{pars}</p:childTnLst></p:cTn>
<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"""

_PAR = """<p:par><p:cTn id="{i0}" fill="hold"><p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
<p:childTnLst><p:par><p:cTn id="{i1}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>
<p:childTnLst><p:par><p:cTn id="{i2}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" nodeType="{nt}">
<p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
<p:set><p:cBhvr><p:cTn id="{i3}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>
<p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="{i4}" dur="400"/>
<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect>
</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>"""

def _build_on_click(s, spids):
    """Fade each shape id in on click, in order. Used sparingly (process maps)."""
    if not spids:
        return
    pars, nid = [], 10
    for k, spid in enumerate(spids):
        pars.append(_PAR.format(i0=nid, i1=nid+1, i2=nid+2, i3=nid+3, i4=nid+4,
                                spid=spid, nt="clickEffect" if k == 0 else "afterEffect"))
        nid += 10
    sld = s._element
    for old in sld.findall(qn("p:timing")):
        sld.remove(old)
    sld.append(etree.fromstring(_TIMING.format(pars="".join(pars))))

def connector(s, x1, y1, x2, y2, color, width=Pt(2.0), arrow=True):
    """A REAL PowerPoint connector line (not a typed arrow glyph)."""
    cx, cy = min(x1, x2), min(y1, y2)
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color; cn.line.width = width
    ln = cn.line._get_or_add_ln()
    if arrow:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle"); tail.set("w", "med"); tail.set("len", "med")
    return cn

def chevron(s, x, y, w, h, color):
    """A staged chevron shape — the real CHEVRON autoshape, not a text glyph."""
    sp = s.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def roundrect(s, x, y, w, h, color, line=None, adj=0.10):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    try: sp.adjustments[0] = adj
    except Exception: pass
    return sp

def diamond(s, x, y, w, h, color, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    return sp

def label_in(sp, text, size, color, bold=True):
    """Put centred text inside an autoshape."""
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Arial"
    return sp

PAGE={"n":1}   # the cover is slide 1 and carries no number, so numbering starts at 2
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}",9,GREY,False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,GREY,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,GREY,False)]],align=PP_ALIGN.RIGHT)
def _ellipsis(text,limit):
    """Truncate on a WORD boundary with an ellipsis — never slice mid-word."""
    t=" ".join(str(text).split())
    if len(t)<=limit: return t
    cut=t[:limit]
    sp=cut.rfind(" ")
    if sp>limit*0.55: cut=cut[:sp]
    return cut.rstrip(" ,.;:-—(") + "…"

def _short_cmd(cmd,limit=30):
    """Shorten a shell command for a caption. Paths/URLs have no spaces, so a
    word-boundary ellipsis cuts them mid-token — drop the runner prefix and keep the
    meaningful tail (the script/target) instead."""
    c=" ".join(str(cmd).split())
    for pre in ("uv run python ","uv run ","python3 ","python ","bash ","sh "):
        if c.startswith(pre): c=c[len(pre):]; break
    if len(c)<=limit: return c
    parts=c.split(" ")
    head=parts[0]
    rest=" ".join(parts[1:]).strip()
    if rest:
        # a URL argument → keep the verb and the host, drop the path
        m=re.match(r'https?://([^/\s]+)',rest)
        if m:
            for cand in (f"{head} {m.group(1)}/…", f"{head} {m.group(1)}"):
                if len(cand)<=limit: return cand
        # "git clone …/repo" — keep the verb plus the final path segment
        tail=rest.rstrip("/").split("/")[-1]
        cand=f"{head} …/{tail}" if "/" in rest else f"{head} {tail}"
        if len(cand)<=limit: return cand
        if len(head)+2<=limit: return _ellipsis(head,limit-2)+" …"
    if "/" in c:                   # a bare path → keep the last segment
        tail=c.rstrip("/").split("/")[-1]
        if len(tail)+2<=limit: return "…/"+tail
    return _ellipsis(c,limit)

def _fit_title(title,size=29):
    """Shrink long titles so they never wrap into the hairline rule below."""
    n=len(title)
    if n<=52: return size
    if n<=66: return 25
    if n<=82: return 22
    return 20

def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),Inches(1.55),kcolor)
    if kicker: txt(s,Inches(0.85),Inches(0.5),Inches(11.6),Inches(0.4),[[(kicker,14,kcolor,True)]])
    txt(s,Inches(0.85),Inches(0.88),Inches(11.9),Inches(0.78),
        [[(title,_fit_title(title),INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
    rect(s,Inches(0.85),Inches(1.7),Inches(11.63),Inches(0.02),LINE)
    return s
def _logo(name):
    p=os.path.join(ASSETS,name)
    return p if os.path.exists(p) else None

# ---------------- slide templates ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    # course badge (top-right) — Gemini Agent ADK badge, else text fallback
    badge=_logo("gemini-adk-badge.png")
    if badge:
        s.shapes.add_picture(badge,Inches(10.35),Inches(0.6),width=Inches(2.2))
    else:
        rect(s,Inches(10.75),Inches(0.66),Inches(1.85),Inches(1.12),BLUE)
        txt(s,Inches(10.75),Inches(0.78),Inches(1.85),Inches(0.52),[[("GEMINI",20,WHITE,True)]],align=PP_ALIGN.CENTER)
        txt(s,Inches(10.75),Inches(1.32),Inches(1.85),Inches(0.4),[[("AGENT ADK",10,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,40,INK,True)]])
    rect(s,Inches(0.92),Inches(4.75),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(5.05),Inches(12),Inches(1.4),
        [[(f"WSQ Course Code: {C.COURSE_CODE}",16,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])

def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    s=head(slide(),title,kicker); bullets(s,Inches(0.85),Inches(1.95),Inches(11.6),Inches(4.9),items,size=size); footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=head(slide(),title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),LIGHT)
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=16)
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=16,mcolor=TEAL); footer(s); return s
def cards3(title,cards,kicker):
    s=head(slide(),title,kicker); xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),Inches(3.65),Inches(4.7),LIGHT); rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),Inches(3.2),Inches(3.4),c[2],size=14,mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
import math
PALETTE=[BLUE,TEAL,VIOLET,AMBER]
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    """Grid of light panels, each with a coloured icon/number badge + text.
    items: list of strings (or (title,caption) tuples). Much richer than a bullet list."""
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,INK,True)],[(it[1],size-2,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def flow_h(title,steps,kicker=None,color=BLUE):
    """Horizontal numbered flow: coloured chips connected by chevrons."""
    s=head(slide(),title,kicker,kcolor=color)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.55); ch=Inches(3.15); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.16),int(y+Inches(1.55)),cw-Inches(0.32),int(ch-Inches(1.7)),[[(st,14,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def process_map(title,stages,kicker=None,color=BLUE,synthesis=None,animate=True):
    """STAGED PROCESS MAP — real rounded-rect stages joined by real connectors with
    arrowheads, each stage numbered, optional synthesis band. stages: list of
    (label, detail). This replaces flow_h wherever a genuine process is taught."""
    s=head(slide(),title,kicker,kcolor=color)
    n=len(stages); X0=Inches(0.85); TOTW=Inches(11.63)
    gap=Inches(0.42); cw=int((TOTW-gap*(n-1))/n)
    y=Inches(2.35); ch=Inches(2.35) if synthesis else Inches(3.0)
    spids=[]
    for i,st in enumerate(stages):
        lbl,detail=(st if isinstance(st,tuple) else (st,""))
        x=int(X0+(cw+gap)*i)
        box=roundrect(s,x,y,cw,ch,LIGHT,line=LINE)
        rect(s,x,y,cw,Inches(0.11),color)
        bd=Inches(0.62)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.34)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.34)),bd,bd,[[(str(i+1),24,WHITE,True)]],
            align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        # label sits in its own fixed band; the caption gets a single line beneath it
        lbl=_ellipsis(lbl,44)
        lsz=13 if len(lbl)<=28 else (12 if len(lbl)<=38 else 11)
        txt(s,x+Inches(0.14),int(y+Inches(1.06)),cw-Inches(0.28),Inches(0.94),
            [[(lbl,lsz,INK,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space=0)
        if detail:
            # ONE line only, sized to the card — a 2-line caption overflows the card bottom
            det=_ellipsis(detail,26)
            txt(s,x+Inches(0.1),int(y+ch-Inches(0.52)),cw-Inches(0.2),Inches(0.38),
                [[(det,9,GREY,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space=0)
        spids.append(box.shape_id)
        if i<n-1:   # REAL connector in the gap, with an arrowhead
            cy=int(y+ch/2)
            cn=connector(s,int(x+cw+Inches(0.06)),cy,int(x+cw+gap-Inches(0.06)),cy,color)
            spids.append(cn.shape_id)
    if synthesis:
        by=int(y+ch+Inches(0.34))
        rect(s,Inches(0.85),by,Inches(11.63),Inches(1.15),LIGHT)
        rect(s,Inches(0.85),by,Inches(0.11),Inches(1.15),color)
        txt(s,Inches(1.15),int(by+Inches(0.12)),Inches(11.1),Inches(0.32),
            [[(synthesis[0].upper(),11,color,True)]])
        txt(s,Inches(1.15),int(by+Inches(0.46)),Inches(11.1),Inches(0.62),
            [[(synthesis[1],13,INK,False)]])
    if animate: _build_on_click(s,spids)
    footer(s); return s

def decision_map(title,question,yes,no,kicker=None,color=VIOLET,note=None):
    """A real decision diamond with two branches drawn as connectors — used for
    'which pattern do I choose' teaching moments."""
    s=head(slide(),title,kicker,kcolor=color)
    # a diamond's usable text area is ~50% of its box — size generously or text spills
    # past the facets. 4.6 x 2.7 fits 2-3 short lines at 12pt.
    dx,dy,dw,dh=Inches(0.85),Inches(2.75),Inches(4.6),Inches(2.7)
    d=diamond(s,dx,dy,dw,dh,color); label_in(d,question,12,WHITE)
    bx=Inches(6.1); bw=Inches(6.35); bh=Inches(1.5)
    ys=[Inches(2.15),Inches(4.35)]
    for (hdr,items),by,col in zip([yes,no],ys,[TEAL,AMBER]):
        b=roundrect(s,bx,by,bw,bh,LIGHT,line=LINE)
        rect(s,bx,by,Inches(0.11),bh,col)
        txt(s,bx+Inches(0.3),int(by+Inches(0.16)),bw-Inches(0.55),Inches(0.4),[[(hdr,15,col,True)]])
        txt(s,bx+Inches(0.3),int(by+Inches(0.6)),bw-Inches(0.55),Inches(0.82),[[(items,12,INK,False)]])
        connector(s,int(dx+dw),int(dy+dh/2),bx,int(by+bh/2),col)
    if note:
        rect(s,Inches(0.85),Inches(6.15),Inches(11.63),Inches(0.72),LIGHT)
        txt(s,Inches(1.15),Inches(6.28),Inches(11.1),Inches(0.5),[[(note,12,GREY,False)]])
    footer(s); return s

def compare_table(title,headers,rows,kicker=None,accent=BLUE,note=None):
    """A real comparison matrix — the substantive alternative to two bullet columns."""
    s=head(slide(),title,kicker,kcolor=accent)
    X0=Inches(0.85); TOTW=Inches(11.63); ncol=len(headers)
    first=int(TOTW*0.26); rest=int((TOTW-first)/(ncol-1))
    widths=[first]+[rest]*(ncol-1)
    y=Inches(1.95); hh=Inches(0.52)
    x=X0
    for i,h in enumerate(headers):
        col=accent if i==0 else PALETTE[(i-1)%len(PALETTE)]
        rect(s,x,y,widths[i],hh,col)
        txt(s,x+Inches(0.14),y,widths[i]-Inches(0.28),hh,[[(h,13,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)
        x+=widths[i]
    # body must end above the note band AND the footer (7.05). Budget it explicitly.
    ry=int(y+hh)
    BOTTOM=Inches(6.88)-(Inches(1.24) if note else Emu(0))
    avail=BOTTOM-ry
    rh=int(min(Inches(0.78),avail/max(len(rows),1)))
    for r,row in enumerate(rows):
        x=X0
        for i,cell in enumerate(row):
            fill=LIGHT if r%2==0 else WHITE
            rect(s,x,ry,widths[i],rh,fill,line=LINE)
            bold=(i==0)
            txt(s,x+Inches(0.14),ry,widths[i]-Inches(0.28),rh,
                [[(cell,11.5,INK if bold else GREY,bold)]],anchor=MSO_ANCHOR.MIDDLE)
            x+=widths[i]
        ry+=rh
    if note:
        rect(s,X0,int(ry+Inches(0.22)),TOTW,Inches(0.92),LIGHT)
        rect(s,X0,int(ry+Inches(0.22)),Inches(0.11),Inches(0.92),accent)
        txt(s,X0+Inches(0.3),int(ry+Inches(0.34)),TOTW-Inches(0.6),Inches(0.66),
            [[("WHEN IT MATTERS  ",11,accent,True),(note,12,INK,False)]])
    footer(s); return s

def worked_example(title,intro,code,explain,kicker=None,accent=TEAL):
    """A worked example: the code on the left, the line-by-line reading on the right.
    This is what turns a decorative lab slide into a teaching slide."""
    s=head(slide(),title,kicker,kcolor=accent)
    txt(s,Inches(0.85),Inches(1.9),Inches(11.63),Inches(0.46),[[(intro,15,GREY,False)]])
    cx,cw=Inches(0.85),Inches(6.5)
    rect(s,cx,Inches(2.5),cw,Inches(4.15),RGBColor(0x0B,0x12,0x20))
    tb=s.shapes.add_textbox(cx+Inches(0.22),Inches(2.62),cw-Inches(0.44),Inches(3.9))
    tf=tb.text_frame; tf.word_wrap=True
    for i,ln in enumerate(code):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(2)
        r=p.add_run(); r.text=ln; r.font.size=Pt(11); r.font.name="Consolas"
        col=RGBColor(0x9C,0xDC,0xFE)
        st=ln.strip()
        if st.startswith("#"): col=RGBColor(0x6A,0x99,0x55)
        elif "=" in ln and not st.startswith(("def","class")): col=RGBColor(0xD4,0xD4,0xD4)
        if st.startswith(("def ","class ","from ","import ")): col=RGBColor(0xC5,0x86,0xC0)
        r.font.color.rgb=col
    ex,ew=Inches(7.65),Inches(4.83)
    for i,(lbl,body) in enumerate(explain[:4]):
        y=int(Inches(2.5)+(Inches(1.0)+Inches(0.05))*i)
        col=PALETTE[i%len(PALETTE)]
        rect(s,ex,y,ew,Inches(1.0),LIGHT); rect(s,ex,y,Inches(0.09),Inches(1.0),col)
        txt(s,ex+Inches(0.26),int(y+Inches(0.1)),ew-Inches(0.45),Inches(0.32),[[(lbl,12,col,True)]])
        txt(s,ex+Inches(0.26),int(y+Inches(0.42)),ew-Inches(0.45),Inches(0.52),[[(body,11,INK,False)]])
    footer(s); return s

def steps_slide(act_title,steps,kicker,accent=TEAL,part=None,start=1):
    """Substantive lab procedure: numbered steps WITH their commands, up to 5 per
    slide. Replaces the old one-sentence step slide."""
    s=head(slide(),act_title+(f" — {part}" if part else ""),kicker,kcolor=accent)
    y0=Inches(1.92); n=len(steps); gapy=Inches(0.1)
    AVAIL=Inches(4.92)          # 1.92 → 6.84, clear of the 7.05 footer
    rh=int(min(Inches(1.12),(AVAIL-gapy*(n-1))/max(n,1)))
    for i,(text,cmd) in enumerate(steps):
        y=int(y0+(rh+gapy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,Inches(0.85),y,Inches(11.63),rh,LIGHT); rect(s,Inches(0.85),y,Inches(0.09),rh,col)
        bd=Inches(0.4)
        oval(s,Inches(1.06),int(y+rh/2-bd/2),bd,bd,col)
        txt(s,Inches(1.06),int(y+rh/2-bd/2),bd,bd,[[(str(start+i),13,WHITE,True)]],
            align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        if cmd:
            txt(s,Inches(1.62),int(y+Inches(0.08)),Inches(10.6),Inches(0.34),[[(text,12.5,INK,True)]])
            rect(s,Inches(1.62),int(y+Inches(0.44)),Inches(10.5),int(rh-Inches(0.54)),RGBColor(0x0B,0x12,0x20))
            one=cmd.split("\n")[0]
            if len(one)>96: one=one[:93]+"..."
            txt(s,Inches(1.78),int(y+Inches(0.44)),Inches(10.2),int(rh-Inches(0.54)),
                [[("$ "+one,10.5,RGBColor(0x9C,0xDC,0xFE),False)]],anchor=MSO_ANCHOR.MIDDLE)
        else:
            txt(s,Inches(1.62),y,Inches(10.6),rh,[[(text,12.5,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s

def chart_slide(title,categories,series,kicker=None,accent=BLUE,
                kind="bar",insight=None,number_format='0'):
    """NATIVE PowerPoint chart (fully editable, not a picture) + an insight band.
    kind: bar | column | line | pie | doughnut. series: list of (name, values)."""
    s=head(slide(),title,kicker,kcolor=accent)
    cd=ChartData(); cd.categories=categories
    for nm,vals in series: cd.add_series(nm,vals,number_format)
    ctype={"bar":XL_CHART_TYPE.BAR_CLUSTERED,"column":XL_CHART_TYPE.COLUMN_CLUSTERED,
           "line":XL_CHART_TYPE.LINE_MARKERS,"pie":XL_CHART_TYPE.PIE,
           "doughnut":XL_CHART_TYPE.DOUGHNUT}.get(kind,XL_CHART_TYPE.COLUMN_CLUSTERED)
    ch_h=Inches(3.55) if insight else Inches(4.75)
    gf=s.shapes.add_chart(ctype,Inches(0.85),Inches(1.95),Inches(11.63),ch_h,cd)
    ch=gf.chart
    ch.has_title=False
    ch.font.size=Pt(12); ch.font.name="Arial"; ch.font.color.rgb=INK
    if kind in ("pie","doughnut") or len(series)>1:
        ch.has_legend=True; ch.legend.position=XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout=False; ch.legend.font.size=Pt(11)
    else:
        ch.has_legend=False
    # house palette per point (pie/doughnut) or per series (bar/column/line)
    try:
        if kind in ("pie","doughnut"):
            pts=ch.plots[0]
            for i,pt in enumerate(pts.points):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb=PALETTE[i%len(PALETTE)]
        else:
            for i,sr in enumerate(ch.series):
                col=PALETTE[i%len(PALETTE)]
                if kind=="line":
                    sr.format.line.color.rgb=col; sr.format.line.width=Pt(2.5)
                else:
                    sr.format.fill.solid(); sr.format.fill.fore_color.rgb=col
    except Exception:
        pass
    try:
        pl=ch.plots[0]; pl.has_data_labels=True
        dl=pl.data_labels; dl.font.size=Pt(10); dl.font.color.rgb=INK
        dl.number_format=number_format; dl.number_format_is_linked=False
        if kind in ("pie","doughnut"): dl.position=XL_LABEL_POSITION.OUTSIDE_END
    except Exception:
        pass
    if insight:
        by=Inches(5.72)
        rect(s,Inches(0.85),by,Inches(11.63),Inches(1.1),LIGHT)
        rect(s,Inches(0.85),by,Inches(0.11),Inches(1.1),accent)
        txt(s,Inches(1.15),int(by+Inches(0.12)),Inches(11.1),Inches(0.3),
            [[("WHAT THE DATA SHOWS",11,accent,True)]])
        txt(s,Inches(1.15),int(by+Inches(0.44)),Inches(11.1),Inches(0.6),
            [[(insight,12,INK,False)]])
    footer(s); return s

def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    """Profile-card layout: avatar badge + name/role panel on the left, labelled
    info tiles on the right. rows: list of (LABEL, value); blank value → fill-in line."""
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),LIGHT); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,GREY,False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,INK,False)] if val else [("____________________________________________",13,LINE,False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def activity_overview(tag,title,desc,build,services,kicker,objective=None,test=None):
    """Lab briefing — now a full teaching slide: the tag chip, the description, and a
    3-tile band covering objective / deliverable / toolchain, plus the success test."""
    s=head(slide(),title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.88),Inches(1.7),Inches(0.46),TEAL)
    txt(s,Inches(0.85),Inches(1.88),Inches(1.7),Inches(0.46),[[(tag,15,WHITE,True)]],
        align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if objective:
        txt(s,Inches(2.72),Inches(1.88),Inches(9.7),Inches(0.46),
            [[(objective,12,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(0.85),Inches(2.5),Inches(11.63),Inches(1.15),[[(desc,17,INK,False)]])
    tiles=[(BLUE,"YOU'LL BUILD",build),(TEAL,"TOOLCHAIN",services),
           (VIOLET,"DONE WHEN",test or "The lab runs end to end without error.")]
    tw=Inches(3.71); xs=[Inches(0.85),Inches(4.81),Inches(8.77)]
    for (col,lbl,body),x in zip(tiles,xs):
        rect(s,x,Inches(3.8),tw,Inches(2.05),LIGHT); rect(s,x,Inches(3.8),tw,Inches(0.1),col)
        txt(s,x+Inches(0.24),Inches(3.98),tw-Inches(0.45),Inches(0.34),[[(lbl,11,col,True)]])
        txt(s,x+Inches(0.24),Inches(4.34),tw-Inches(0.45),Inches(1.4),[[(body,12,INK,False)]])
    footer(s); return s
def step_slide(kicker,act_title,n,total,text,cmd=""):
    s=head(slide(),act_title,kicker,TEAL)
    oval(s,Inches(0.85),Inches(2.5),Inches(1.4),Inches(1.4),TEAL)
    txt(s,Inches(0.85),Inches(2.74),Inches(1.4),Inches(0.9),[[(str(n),38,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.95),Inches(1.95),Inches(11),Inches(0.4),[[(f"STEP {n} OF {total}",13,GREY,True)]])
    txt(s,Inches(2.55),Inches(2.4),Inches(10.1),Inches(1.3),[[(text,23,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if cmd:
        rect(s,Inches(2.55),Inches(4.15),Inches(10.1),Inches(0.95),RGBColor(0x0B,0x12,0x20))
        txt(s,Inches(2.8),Inches(4.28),Inches(9.7),Inches(0.7),[[("$ "+cmd,13,RGBColor(0x9C,0xDC,0xFE),False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def test_slide(act_title,text,kicker,troubleshoot=None):
    """Verification slide — the success criterion PLUS a troubleshooting band, so it
    teaches diagnosis rather than stating one sentence."""
    s=head(slide(),act_title,kicker,TEAL)
    GREEN=RGBColor(0x12,0x7A,0x3E)
    rect(s,Inches(0.85),Inches(1.95),Inches(11.63),Inches(2.15),RGBColor(0xE8,0xF7,0xEE))
    rect(s,Inches(0.85),Inches(1.95),Inches(0.11),Inches(2.15),GREEN)
    txt(s,Inches(1.2),Inches(2.12),Inches(11),Inches(0.44),[[("✅  Expected result",17,GREEN,True)]])
    txt(s,Inches(1.2),Inches(2.62),Inches(11.0),Inches(1.35),[[(text,15,INK,False)]])
    tb=troubleshoot or [
        ("Nothing happens","Check the .env file is in the labs folder and the key has no quotes or spaces."),
        ("Auth or 401 error","Re-copy the API key from AI Studio; confirm GOOGLE_GENAI_USE_VERTEXAI=0."),
        ("ModuleNotFoundError","Run uv sync again, and prefix commands with uv run so the venv is used."),
    ]
    txt(s,Inches(0.85),Inches(4.32),Inches(11.63),Inches(0.34),
        [[("IF IT DOESN'T WORK",11,AMBER,True)]])
    tw=Inches(3.71); xs=[Inches(0.85),Inches(4.81),Inches(8.77)]
    for i,(sym,fix) in enumerate(tb[:3]):
        x=xs[i]
        rect(s,x,Inches(4.7),tw,Inches(1.72),LIGHT); rect(s,x,Inches(4.7),tw,Inches(0.09),AMBER)
        txt(s,x+Inches(0.24),Inches(4.87),tw-Inches(0.45),Inches(0.36),[[(sym,12.5,INK,True)]])
        txt(s,x+Inches(0.24),Inches(5.26),tw-Inches(0.45),Inches(1.05),[[(fix,11,GREY,False)]])
    footer(s); return s
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1

# ============================================================ BUILD
cover()

# ---------------- ADMIN ----------------
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
tile_grid("Digital Attendance (Mandatory)",[
 ("Three times a day","Take the AM, PM and Assessment digital attendance — mandatory for every WSQ-funded course."),
 ("Trainer shows the QR","The trainer or administrator displays the digital attendance QR code from the SSG portal."),
 ("Scan and submit","Scan the QR code with your mobile phone camera and submit your attendance."),
 ("75% minimum","A minimum of 75% attendance is required to be eligible for assessment and funding.")],
 kicker="TRAQOM · SSG DIGITAL ATTENDANCE",cols=2,size=15)
trainer_slide("YOUR TRAINER · GENERAL","Your Trainer","General Trainer template —\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",C.TRAINER,"Principal Trainer\nTertiary Infotech Academy Pte Ltd",
 [("Role","Principal Trainer, Tertiary Infotech Academy Pte Ltd"),
  ("Expertise","Generative AI, Large Language Models, agentic AI systems and applied machine learning."),
  ("Delivers","WSQ courses on AI, LLM application development, prompt engineering and data science."),
  ("Founder","Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
 initials="AA",accent=BLUE)
content("Let's Know Each Other",[
 "Your name, organisation and role.",
 "Your experience with Python, APIs and Generative AI tools (if any).",
 "What you would like to automate or build with AI agents after this course."],kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Participate actively — no question is too small.",
 "Mutual respect: agree to disagree.","One conversation at a time.",
 "Be punctual; return from breaks on time.","75% attendance is required."],
 kicker="HOUSEKEEPING",cols=2,size=15)
tile_grid("Download Your Course Material",[
 ("1 · Go to the LMS portal","Open https://lms-tms.tertiaryinfotech.com in your browser."),
 ("2 · Log in","Sign in with the account e-mail you used to register for this course."),
 ("3 · Open this course","Select 'Develop Multi AI Agent Applications with Gemini Agent ADK'."),
 ("4 · Download","Download the Trainer Slides, Learner Guide and Lesson Plan (PDF)."),
 ("5 · Get the labs","Clone the lab repository from GitHub — the link is on the course page."),
 ("6 · Keep them open","You may use these materials during the open-book assessment.")],
 kicker="LMS / TMS  ·  lms-tms.tertiaryinfotech.com",cols=2,size=14)
tile_grid("Skills Framework Alignment",[
 ("TSC Title", C.TSC_TITLE),
 ("TSC Code", C.TSC_CODE),
 ("A1 / A2", "Analyse algorithms in AI applications; correlate algorithm design with efficiency."),
 ("A3 / A4", "Identify and compare the strengths and limitations of AI applications."),
 ("A5", "Assess the feasibility of AI applications for engineering processes."),
 ("A6", "Assess the improvements AI applications bring to engineering processes.")],
 kicker="SKILLS FRAMEWORK  ·  TSC",cols=2,size=14,accent=VIOLET)
two_col("Lesson Plan — 2 Days, 8 hours/day",[
 (f"Day 1 — {C.DAY_THEMES[1]}",0),
 ("Digital Attendance (AM) · Introductions · Learning Outcomes",1),
 ("Topic 1: Overview of Agentic AI in Gemini ADK (Labs 1–4)",1),
 ("Lunch Break · Digital Attendance (PM)",1),
 ("Topic 2: Build A Multi Agent App with Gemini ADK begins (Labs 5–8)",1)],
 [(f"Day 2 — {C.DAY_THEMES[2]}",0),
 ("Topic 2 continues (Labs 9–12)",1),
 ("Topic 3: Build Agentic AI RAG in Gemini ADK (Labs 13–15)",1),
 ("Topic 4: Agentic AI App with Streamlit (Labs 16–18)",1),
 ("TRAQOM Survey · Digital Attendance (Assessment)",1),
 ("Final Assessment (WA + PP)",1)],
 kicker="SCHEDULE",lhead="Day 1",rhead="Day 2")
tile_grid("Learning Outcomes",[
 ("LO1 · Analyse LLM applications","Analyze the range of LLM applications using Generative AI and identify their industrial use cases."),
 ("LO2 · Establish Gemini designs","Establish Google Gemini GAI designs and assess improvements on engineering processes."),
 ("LO3 · Develop LLM applications","Develop LLM applications with the Agent Development Kit and assess their feasibility."),
 ("LO4 · Evaluate RAG","Evaluate the performance effectiveness of Retrieval Augmented Generation.")],
 kicker="WHAT YOU'LL ACHIEVE",cols=2,size=15)
tile_grid("Course Outline",[
 ("Topic 1 — Overview of Agentic AI in Gemini ADK","LLM & agentic AI foundations · the Gemini model family · ADK architecture · your first agent · tools."),
 ("Topic 2 — Build A Multi Agent App with Gemini ADK","Sessions & state · handoff · sub-agents · workflow agents · guardrails · structured output · MCP."),
 ("Topic 3 — Build Agentic AI RAG in Gemini ADK","The RAG pipeline · loading & splitting · embeddings · vector stores · retrieval · evaluation."),
 ("Topic 4 — Agentic AI App with Gemini ADK and Streamlit","YAML agents · Streamlit chat UI · async Runner · deployment and feasibility · capstone.")],
 kicker="FOUR TOPICS  ·  18 HANDS-ON LABS",cols=1,size=15)
tile_grid("Briefing for Assessment",[
 ("Do · Clear your desk","Place phones and other materials under the table or on the floor."),
 ("Don't · No recording","No photos or recording of assessment scripts."),
 ("Don't · No discussion","Work individually — no discussion during the assessment."),
 ("Do · Black or blue pen","Use a black or blue pen for hard-copy assessments."),
 ("Don't · No correction fluid","No liquid paper or correction tape may be used."),
 ("Do · Stop on time","Scripts are collected when time is up.")],
 kicker="BEFORE YOU START",cols=2,size=14,accent=AMBER)
tile_grid("Assessment",[
 ("Written Assessment (WA)","Short-Answer Questions (SAQ) · 1 hour · open book. Tests the underpinning knowledge taught in the slides."),
 ("Practical Performance (PP)","Hands-on Gemini ADK agent build tasks · 1 hour · open book. Tests the ability you built in the labs."),
 ("Open book","You may use the course slides, the Learner Guide and approved materials only."),
 ("Eligibility","A minimum of 75% attendance is required to be eligible for assessment and funding."),
 ("Result","You are assessed as Competent (C) or Not Yet Competent (NYC) on each instrument."),
 ("Appeals","An appeal process is available if you wish to contest an assessment outcome.")],
 kicker="FINAL ASSESSMENT",cols=2,size=14)
process_map("Assessment Flow",[
 ("TRAQOM survey","Scan the QR code on the LMS"),
 ("Digital attendance","Scan the SSG assessment QR"),
 ("Sit WA then PP","Open book · 1 hour each"),
 ("Submit on the LMS","Upload your completed answers"),
 ("Sign the record","Sign the Assessment Summary Record")],
 kicker="ON ASSESSMENT DAY",color=BLUE,
 synthesis=("REMEMBER","All five steps are mandatory for WSQ funding — missing the digital attendance or the TRAQOM survey can invalidate your claim."))
tile_grid("Criteria for Funding",[
 ("Attendance","A minimum attendance rate of 75%, based on the SSG Digital Attendance record."),
 ("Assessment","Complete both assessment components and be assessed as 'Competent'."),
 ("Digital attendance","Scan the SSG QR code for AM, PM and Assessment on every training day."),
 ("TRAQOM survey","Complete the mandatory TRAQOM course feedback survey on the LMS.")],
 kicker="WSQ FUNDING",cols=2,size=15,accent=AMBER)
tile_grid("Set Up Before We Start",[
 ("Google account","Sign in to Google AI Studio at https://aistudio.google.com to create a free API key."),
 ("Python 3.13","Install Python 3.13 and the uv package manager for the lab environment."),
 ("Lab repository","Clone the course labs from GitHub — the link is on your LMS course page."),
 ("Editor","Use VS Code or any editor you are comfortable with for editing agent.py files."),
 ("Optional keys","OpenWeather and Tavily free API keys are used by the tool-calling labs."),
 ("ADK docs","Bookmark the official ADK documentation at https://google.github.io/adk-docs/.")],
 kicker="LAB ENVIRONMENT",cols=2,size=14,accent=TEAL)

# ---------------- CORE CONCEPTS ----------------
section("CORE CONCEPTS","From Large Language Models to AI Agents","")
tile_grid("What is a Large Language Model?",[
 ("A model of language","An LLM such as Gemini, GPT-4 or Claude is trained to understand and generate human language."),
 ("Why 'large'","Large in both the volume of training data and the model's complexity — hundreds of billions of parameters."),
 ("Autoregressive","It predicts the next token, one token at a time, conditioned on everything written so far."),
 ("Trained then inferred","Training builds the model from a huge corpus; inference is the prompt-and-response you use every day.")],
 kicker="FOUNDATIONS",cols=2,size=15)
flow_h("How an LLM Produces an Answer",[
 "Your prompt is broken into tokens",
 "The model attends over the full context",
 "It predicts the most likely next token",
 "The token is appended and the loop repeats",
 "Generation stops at a stop token or limit"],kicker="AUTOREGRESSIVE GENERATION")
tile_grid("Opportunities of LLM Applications",[
 "Content creation and writing assistance","Customer support and chatbots",
 "Language translation and localisation","Educational and tutoring tools",
 "Business intelligence and analytics","Coding and software development",
 "Legal and compliance assistance","Healthcare and clinical documentation support",
 "Enhanced search and knowledge discovery","Accessibility for persons with disabilities"],
 kicker="LO1 · WHERE LLMs ARE USED",cols=2,size=14)
cards3("Industry Use Cases of Generative AI",[
 (BLUE,"Finance & Legal",["Automated financial reporting","Risk assessment and analysis","Fraud detection","Contract review and drafting","Regulatory compliance monitoring"]),
 (TEAL,"Marketing & eCommerce",["Content and copy generation","Personalised recommendations","Customer sentiment analysis","SEO and campaign optimisation","Product description generation"]),
 (VIOLET,"Education & Healthcare",["Personalised learning assistance","Automated grading and feedback","Clinical documentation","Medical imaging support","Virtual health assistants"])],
 kicker="LO1 · INDUSTRIAL USE CASES")
tile_grid("Popular Large Language Models",[
 ("Google Gemini","Natively multimodal; Ultra, Pro, Flash and Nano variants — the model family used in this course."),
 ("OpenAI GPT","Widely adopted general-purpose family; reachable from ADK through the LiteLlm wrapper."),
 ("Anthropic Claude","Strong long-context reasoning and writing; also reachable through LiteLlm."),
 ("Open models","Meta Llama and Mistral can be self-hosted or run locally through Ollama for data-residency needs.")],
 kicker="THE LANDSCAPE",cols=2,size=15)
big_statement("An LLM answers. An agent acts.","Agentic AI adds reasoning, memory, tool use and autonomy on top of a language model — so it can complete tasks, not just produce text.","WHY AGENTS",color=BLUE)
tile_grid("What Makes a System 'Agentic'?",[
 ("Reasoning","The model plans a sequence of steps instead of answering in one shot."),
 ("Tool use","It calls functions and APIs to fetch live data or take real actions."),
 ("Memory","Session state lets it carry context across many turns of a conversation."),
 ("Autonomy","It decides which tool or which specialist agent handles the request."),
 ("Collaboration","Multiple specialised agents divide the work and hand off between themselves."),
 ("Guardrails","Callbacks and policies constrain what the agent is allowed to do.")],
 kicker="THE SIX INGREDIENTS",cols=2,size=14,accent=VIOLET)
compare_table("LLM vs Agent vs Multi-Agent System",
 ["Dimension","Plain LLM","Single Agent","Multi-Agent System"],
 [["What it does","Generates text","Completes a task","Divides work across specialists"],
  ["Memory","None — stateless","Session state","Shared session + per-agent state"],
  ["Tools","None","Its own tool set","Each agent has its own tools"],
  ["Best for","Drafting, summarising","One well-defined job","Broad domains, varied requests"],
  ["Main risk","Hallucination","Overloaded instruction","Wrong routing / handoff loops"]],
 kicker="LO1 · CHOOSING THE RIGHT ARCHITECTURE",accent=BLUE,
 note="Add agents only when one instruction can no longer describe the job clearly — complexity you don't need is complexity you must still debug.")
process_map("The Agentic Loop — Reason, Act, Observe",
 [("Perceive","Read the user request and session state"),
  ("Reason","Plan which step or tool comes next"),
  ("Act","Call a tool, or hand off to a sub-agent"),
  ("Observe","Read the tool result back into context"),
  ("Respond","Answer, or loop again if not done")],
 kicker="THE PATTERN BEHIND EVERY ADK AGENT",color=VIOLET,
 synthesis=("WHY IT MATTERS","The loop is what separates an agent from a chatbot: it can take several actions before answering, and it decides how many."))
worked_example("Reading an ADK Agent Definition",
 "Every agent in this course is built from these same fields — learn to read them once.",
 ["from google.adk.agents import Agent","","agent = Agent(",
  "    model='gemini-2.0-flash',","    name='banking_assistant',",
  "    description='Handles retail banking queries',",
  "    instruction=(","        'You are a bank assistant. '",
  "        'Never reveal a PIN or OTP.'","    ),","    tools=[get_balance],",")"],
 [("model","Which LLM runs the agent. Flash is fast and cheap; Pro reasons harder."),
  ("name","The identifier used in traces and in agent-to-agent transfer."),
  ("description","Read by a PARENT agent to decide whether to hand off here."),
  ("instruction","The system prompt — role, rules and refusals. Most behaviour lives here.")],
 kicker="LO3 · WORKED EXAMPLE",accent=TEAL)
decision_map("Which Agent Pattern Should You Use?",
 "Can one clear\ninstruction describe\nthe whole job?",
 ("YES — use a single agent","Give it the tools it needs. Simpler to build, cheaper to run and far easier to debug. Most business tasks stop here."),
 ("NO — use sub-agents / workflow","Split by specialism. Use sub-agents for routing by intent, or a SequentialAgent when the order of steps is fixed."),
 kicker="LO3 · DESIGN DECISION",
 note="A common mistake is starting multi-agent. Start single, and split only when the instruction starts listing unrelated responsibilities.")
chart_slide("Where Course Time Goes — Topic Weighting",
 ["Topic 1 — Agentic AI Overview","Topic 2 — Multi-Agent Apps",
  "Topic 3 — Agentic RAG","Topic 4 — Streamlit App"],
 [("Share of course (%)",[int(str(t["weighting"]).replace("%","") or 0) for t in C.TOPICS])],
 kicker="LESSON PLAN · TIME ALLOCATION",accent=BLUE,kind="column",
 insight="Topic 2 carries the largest weighting because multi-agent orchestration is where most of the assessable practical skill sits — plan your revision accordingly.",
 number_format='0"%"')
chart_slide("Cost and Latency by Model Tier",
 ["Flash","Pro","Ultra"],
 [("Relative cost per 1K calls",[1,7,21]),("Relative latency",[1,3,6])],
 kicker="LO2 · ENGINEERING TRADE-OFF",accent=TEAL,kind="column",
 insight="Moving from Flash to Ultra costs roughly 20x more and runs about 6x slower. Reach for the larger model only where the reasoning genuinely requires it — indicative figures for teaching the trade-off, not a price list.")
compare_table("Gemini Model Selection — Trade-offs",
 ["Model","Speed","Cost","Use it when"],
 [["Gemini Flash","Fastest","Lowest","High-volume agents, tool routing, most labs in this course"],
  ["Gemini Pro","Moderate","Medium","Complex reasoning, long documents, quality-critical answers"],
  ["Gemini Ultra","Slowest","Highest","The hardest reasoning tasks where quality outweighs cost"],
  ["Gemini Nano","On-device","None","Offline or mobile scenarios with privacy constraints"]],
 kicker="LO2 · ASSESSING FEASIBILITY",accent=TEAL,
 note="Model choice is an engineering trade-off, not a preference — justify it on latency, cost per call and required reasoning depth.")
two_col("Single Agent vs Multi-Agent",[
 ("Single agent",0),("One instruction, one set of tools",1),
 ("Simple to build and to debug",1),("Instruction grows unwieldy as scope widens",1),
 ("Best for a narrow, well-defined job",1)],
 [("Multi-agent",0),("Each specialist has its own instruction and tools",1),
 ("Scales to broad problem domains",1),("Routing and handoff must be designed carefully",1),
 ("Best for varied requests across several sub-domains",1)],
 kicker="ARCHITECTURE CHOICE",lhead="Simple",rhead="Scalable")
tile_grid("What is Gemini?",[
 ("Built by Google DeepMind","Gemini is Google's flagship model family, built from the ground up for multimodality."),
 ("Natively multimodal","It reasons across text, images, video, audio and code in a single model."),
 ("Long context","Very large context windows allow whole documents, codebases and videos in one prompt."),
 ("Tool-ready","Native function calling makes Gemini a natural fit for agentic applications.")],
 kicker="THE MODEL BEHIND ADK",cols=2,size=15)
cards3("The Gemini Model Family",[
 (BLUE,"Pro",["The mid-tier, balanced option","Strong reasoning, broad multimodal skill","Optimised for performance, cost and latency","The default choice for most agents"]),
 (TEAL,"Flash",["Lightweight and fast","Optimised for speed and efficiency","Low cost per call at high volume","Used in most labs in this course"]),
 (VIOLET,"Ultra & Nano",["Ultra: the most capable, for highly complex tasks","Nano: the most efficient, runs on-device","Nano targets mobile and offline scenarios","Choose by task complexity and where it runs"])],
 kicker="CHOOSING A MODEL")
tile_grid("Model Parameters You Can Tune",[
 ("Temperature","Controls randomness. Low values give focused, repeatable answers; high values give creative, varied ones."),
 ("Top-K","Restricts sampling to the K most probable next tokens. Top-K of 1 always picks the single most likely token."),
 ("Top-P","Samples from the smallest set of tokens whose probabilities sum to P — a dynamic alternative to Top-K."),
 ("Max output tokens","Caps the length of the response, which directly caps cost."),
 ("Stop sequences","Strings that end generation immediately when produced — useful for structured formats."),
 ("Safety settings","Thresholds that filter harmful content across several harm categories.")],
 kicker="LO2 · GEMINI GAI DESIGN",cols=2,size=14)
tile_grid("What is the Agent Development Kit (ADK)?",[
 ("Open-source framework","Google's Python framework for building, evaluating and deploying AI agents."),
 ("Model-agnostic","Runs on Gemini natively, and on OpenAI, Anthropic or local models through LiteLlm."),
 ("Batteries included","Sessions, tools, callbacks, multi-agent orchestration and evaluation are all built in."),
 ("Developer UI","adk web gives you a local browser IDE that shows every event, tool call and trace.")],
 kicker="TOPIC 1 PREVIEW",cols=2,size=15,accent=TEAL)
tile_grid("Anatomy of an ADK Agent",[
 ("model","Which LLM powers the agent — for example gemini-2.0-flash."),
 ("name","A unique identifier used in traces and in agent-to-agent transfers."),
 ("description","What this agent is for. A parent agent reads this to decide whether to hand off to it."),
 ("instruction","The system prompt: the agent's role, rules, tone and boundaries."),
 ("tools","Python functions the agent may call to fetch data or take action."),
 ("sub_agents","Specialist agents this agent may transfer control to.")],
 kicker="THE SIX FIELDS",cols=2,size=14,accent=TEAL)
flow_h("The ADK Agent Loop",[
 "User message enters the Runner",
 "Session history is loaded as context",
 "The model reasons and may emit a tool call",
 "The tool runs and its result returns to the model",
 "The final response is streamed back and stored"],kicker="HOW A TURN EXECUTES",color=TEAL)

# ---------------- TOPICS + ACTIVITIES ----------------
TOPIC_ACTS = {t["num"]: [a for a in ACTIVITIES if a["topic"]==t["num"]] for t in C.TOPICS}
CARD_COLORS=[BLUE,TEAL,VIOLET]
for t in C.TOPICS:
    section(f"TOPIC {t['code']}", t["title"], t["code"], t["subtitle"])
    cons=t["concepts"]
    half=(len(cons)+1)//2
    for ci,chunk in enumerate([cons[:half],cons[half:]]):
        if not chunk: continue
        suffix="" if ci==0 else " (continued)"
        tile_grid(f"Key Concepts — {t['title']}{suffix}", chunk,
                  kicker=f"TOPIC WEIGHTING {t['weighting']}", cols=2, size=14)
    acts=TOPIC_ACTS[t["num"]]
    third=(len(acts)+2)//3
    groups=[acts[i:i+third] for i in range(0,len(acts),third)][:3]
    while len(groups)<3: groups.append([])
    cards=[]
    for gi,g in enumerate(groups):
        _lbl=("—" if not g else (f"Lab {g[0]['num']}" if g[0]['num']==g[-1]['num']
                                 else f"Labs {g[0]['num']}–{g[-1]['num']}"))
        cards.append((CARD_COLORS[gi], _lbl,
                      [a["title"] for a in g] if g else ["—"]))
    cards3(f"Hands-On Labs — {t['title']}", cards, kicker="WHAT YOU'LL DO")
    for a in acts:
        activity_overview(f"LAB {a['num']}", a["title"], a["desc"], a["build"], a["services"],
                          kicker=f"TOPIC {t['code']} · HANDS-ON",
                          objective=a.get("objective"), test=a.get("test"))
        # the lab as a real staged process map (first 5 steps, abridged to stage labels)
        stg=[]
        for (stext,scmd) in a["steps"][:5]:
            lbl=stext.split(",")[0].split(" and ")[0].split(" (")[0].strip()
            cap=scmd.split("\n")[0].strip() if scmd else ""
            stg.append((_ellipsis(lbl,44), _short_cmd(cap,30) if cap else ""))
        if len(stg)>=3:
            process_map(f"How Lab {a['num']} Runs", stg,
                        kicker=f"LAB {a['num']} · PROCESS MAP", color=TEAL,
                        synthesis=("THE POINT", a["objective"]))
        # full procedure, with the actual commands, split across slides of 5
        allsteps=a["steps"]
        chunks=[allsteps[i:i+4] for i in range(0,len(allsteps),4)]
        for ci,ch in enumerate(chunks):
            steps_slide(a["title"], ch, kicker=f"LAB {a['num']} · PROCEDURE",
                        part=(f"Part {ci+1}/{len(chunks)}" if len(chunks)>1 else None),
                        start=ci*4+1)
        test_slide(a["title"], a["test"], kicker=f"LAB {a['num']} · VERIFY")
    content(f"Recap — {t['title']}",
            ["You can now: "+a["objective"] for a in {x["objective"]:x for x in acts}.values()][:6],
            kicker="TOPIC RECAP", size=17)

# ---------------- CLOSE ----------------
section("WRAP-UP","Course Summary & Next Steps","")
tile_grid("What You Achieved",[
 ("LO1 · Analysed LLM applications","Explained how LLMs work and identified Generative AI use cases across finance, marketing, education, healthcare and legal."),
 ("LO2 · Established Gemini designs","Chose Gemini models and tuned parameters, and assessed the improvements agents bring to engineering processes."),
 ("LO3 · Developed LLM applications","Built single agents, tool-using agents, multi-agent handoff, sequential workflows, guardrails, structured output and MCP integrations."),
 ("LO4 · Evaluated RAG","Built an agentic RAG pipeline over your own PDFs and measured retrieval hit rate, groundedness and latency.")],
 kicker="LEARNING OUTCOMES",cols=2,size=14)
tile_grid("Continue Your Learning",[
 ("Official ADK docs","https://google.github.io/adk-docs/ — the reference for every ADK class and pattern."),
 ("Google AI Studio","https://aistudio.google.com — prototype prompts and manage your API keys."),
 ("Rebuild the labs","Redo each lab from a blank file until the agent pattern is automatic."),
 ("Extend the capstone","Add RAG, a guardrail and a Streamlit UI to your own multi-agent application."),
 ("Deploy it","Take your agent to Cloud Run or Vertex AI Agent Engine for production use."),
 ("Share it","Publish your agent to GitHub and gather feedback from real users.")],
 kicker="NEXT STEPS",cols=2,size=14)
tile_grid("Recommended Courses",[(rc,"") for rc in C.RECOMMENDED_COURSES],
 kicker="CONTINUE WITH TERTIARY INFOTECH",cols=1,size=15)
content("Support",[
 "If you have any enquiries during or after the class, you can contact us below.",
 "Email: enquiry@tertiaryinfotech.com","Tel: +65 6100 0613",
 "Website: www.tertiarycourses.com.sg","LMS / TMS: https://lms-tms.tertiaryinfotech.com"],kicker="WE'RE HERE TO HELP")
tile_grid("Assessment",[
 ("Written Assessment (SAQ)","1 hour · open book · short-answer knowledge questions."),
 ("Practical Performance (PP)","1 hour · open book · hands-on Gemini ADK agent build tasks."),
 ("Digital attendance","Remember to take the Assessment digital attendance (TRAQOM) before you start."),
 ("Submit on the LMS","Upload your completed answers at https://lms-tms.tertiaryinfotech.com/")],
 kicker="WRAP-UP",cols=2,size=15)
process_map("Assessment Flow",[
 ("TRAQOM survey","Scan the QR code on the LMS"),
 ("Digital attendance","Scan the SSG assessment QR"),
 ("Sit WA then PP","Open book · 1 hour each"),
 ("Submit on the LMS","Upload your completed answers"),
 ("Sign the record","Sign the Assessment Summary Record")],
 kicker="ON ASSESSMENT DAY",color=BLUE,
 synthesis=("REMEMBER","All five steps are mandatory for WSQ funding — missing the digital attendance or the TRAQOM survey can invalidate your claim."))
tile_grid("Digital Attendance (Mandatory)",[
 ("Three times a day","Take the AM, PM and Assessment digital attendance — mandatory for every WSQ-funded course."),
 ("Trainer shows the QR","The trainer or administrator displays the digital attendance QR code from the SSG portal."),
 ("Scan and submit","Scan the QR code with your mobile phone camera and submit your attendance."),
 ("75% minimum","A minimum of 75% attendance is required to be eligible for assessment and funding.")],
 kicker="TRAQOM · SSG DIGITAL ATTENDANCE",cols=2,size=15)
big_statement("Thank You!","You can now design, build and evaluate multi-agent AI applications with the Gemini Agent Development Kit.","HAPPY BUILDING",color=TEAL)

# ---------------- restrained motion pass ----------------
# ONE family for the whole deck: content fades, dividers push. Applied last so every
# slide is covered exactly once. Nothing spins, flies or bounces.
DIVIDER_IDS=set()
for s in prs.slides:
    txts=[sh.text_frame.text for sh in s.shapes if sh.has_text_frame]
    joined=" ".join(txts)
    is_div=any(k in joined for k in ("COURSE ADMINISTRATION","CORE CONCEPTS","WRAP-UP")) or \
           any(t["title"] in joined and f"TOPIC {t['code']}" in joined for t in C.TOPICS)
    _transition(s,"push" if is_div else "fade",speed="fast" if not is_div else "med")

OUT=os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
prs.save(OUT)
print(f"Saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
